import os
from datetime import datetime, timedelta
from pathlib import Path
import csv
import io
import time
import re
import asyncio
from typing import List
import zipfile
import openpyxl
from openpyxl import Workbook
import uuid
import logging
import img2pdf
from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    ReplyKeyboardMarkup,
    KeyboardButton,
    Bot,
    InputFile
)
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    CallbackContext,
    ContextTypes,
    filters
)
from telegram.constants import ParseMode
from docx2pdf import convert  # DOCX to PDF uchun
from PIL import Image, ImageDraw, ImageFont
from fpdf import FPDF
from docx import Document
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PIL import Image
import tempfile
from io import BytesIO
import subprocess
import sys
import docx
from openpyxl import load_workbook
from pptx import Presentation
from pdf2docx import Converter
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import pytesseract
import PyPDF2
import platform
from contextlib import contextmanager
import shutil
import psutil
import qrcode
import cv2
import numpy as np
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from pptx.util import Inches
from flask import Flask, request, jsonify
import gunicorn

# Tesseract yo‘lini Linux uchun sozlash
pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

# Logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('bot.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

TOKEN = os.environ.get('TOKEN', "7797976277:AAGeRUw7sqMh_PQrPNsISTHs_9cSrXyzFiQ")
ADMIN_ID = int(os.environ.get('ADMIN_ID', 1223308504))

# Global o'zgaruvchilar (oldingi kodingizdan saqlangan)
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB for DOCX to PDF
MAX_EXCEL_SIZE = 5 * 1024 * 1024  # 5MB
ZIP_BUTTON = "zip_button"
CLEAR_BUTTON = "clear_button"
OCR_BUTTON = "ocr_button"
COMPRESS_BUTTON = "compress_button"
COMPRESS_HIGH = "compress_high"
COMPRESS_MEDIUM = "compress_medium"
COMPRESS_MAX = "compress_max"
OCR_MESSAGE = "ocr_message"
OCR_TXT = "ocr_txt"
OCR_DOCX = "ocr_docx"

# Global o'zgaruvchilar Obunaniki
users = {}  # user_id: {'first_name': str, 'username': str}
user_activity = {}
subscription_channels = []
user_subscriptions = {}
pending_messages = {}
check_tasks = {}
confirmation_messages = {}

# Translation dictionaries (saqlangan)
CYRILLIC_TO_LATIN = {
    'а': 'a', 'б': 'b', 'в': 'v', 'г': 'g', 'д': 'd', 'е': 'e', 'ё': 'yo',
    'ж': 'j', 'з': 'z', 'и': 'i', 'й': 'y', 'к': 'k', 'л': 'l', 'м': 'm',
    'н': 'n', 'о': 'o', 'п': 'p', 'р': 'r', 'с': 's', 'т': 't', 'у': 'u',
    'ф': 'f', 'х': 'x', 'ц': 'ts', 'ч': 'ch', 'ш': 'sh', 'щ': 'shch',
    'ъ': "'", 'ы': 'i', 'ь': "'", 'э': 'e', 'ю': 'yu', 'я': 'ya',
    'ў': "o'", 'қ': 'q', 'ғ': "g'", 'ҳ': 'h',
    'А': 'A', 'Б': 'B', 'В': 'V', 'Г': 'G', 'Д': 'D', 'Е': 'E', 'Ё': 'Yo',
    'Ж': 'J', 'З': 'Z', 'И': 'I', 'Й': 'Y', 'К': 'K', 'Л': 'L', 'М': 'M',
    'Н': 'N', 'О': 'O', 'П': 'P', 'Р': 'R', 'С': 'S', 'Т': 'T', 'У': 'U',
    'Ф': 'F', 'Х': 'X', 'Ц': 'Ts', 'Ч': 'Ch', 'Ш': 'Sh', 'Щ': 'Shch',
    'Ъ': "'", 'Ы': 'I', 'Ь': "'", 'Э': 'E', 'Ю': 'Yu', 'Я': 'Ya',
    'Ў': "O'", 'Қ': 'Q', 'Ғ': "G'", 'Ҳ': 'H'
}

LATIN_TO_CYRILLIC = {
    'a': 'а', 'b': 'б', 'd': 'д', 'e': 'е', 'f': 'ф', 'g': 'г',
    'h': 'ҳ', 'i': 'и', 'j': 'ж', 'k': 'к', 'l': 'л', 'm': 'м',
    'n': 'н', 'o': 'о', 'p': 'п', 'q': 'қ', 'r': 'р', 's': 'с',
    't': 'т', 'u': 'у', 'v': 'в', 'x': 'х', 'y': 'й', 'z': 'з',
    "'": 'ъ', '`': 'ъ', 'ʻ': 'ъ', 'ʼ': 'ъ', '’': 'ъ',
    'sh': 'ш', 'ch': 'ч', 'ye': 'е', 'ya': 'я', 'yu': 'ю',
    "g'": 'ғ', "o'": 'ў',
    'A': 'А', 'B': 'Б', 'D': 'Д', 'E': 'Е', 'F': 'Ф', 'G': 'Г',
    'H': 'Ҳ', 'I': 'И', 'J': 'Ж', 'K': 'К', 'L': 'Л', 'M': 'М',
    'N': 'Н', 'O': 'О', 'P': 'П', 'Q': 'Қ', 'R': 'Р', 'S': 'С',
    'T': 'Т', 'U': 'У', 'V': 'В', 'X': 'Х', 'Y': 'Й', 'Z': 'З',
    'Sh': 'Ш', 'Ch': 'Ч', 'Ye': 'Е', 'Ya': 'Я', 'Yu': 'Ю',
    "G'": 'Ғ', "O'": 'Ў', "O‘": 'Ў',
    'SH': 'Ш', 'CH': 'Ч'
}

# User data storage
user_data = {}

class UserData:
    def __init__(self):
        self.files = []
        self.status_message_id = None
        self.images = []
        self.waiting_for_images = False
        self.active_module = None
        self.lock = asyncio.Lock()
        self.ocr_text = None
        self.output_format = None
        self.compressed_sizes = {}
        self.temp_file_path = None
        self.temp_file_name = None
        self.temp_dir = None

@contextmanager
def temp_file_manager():
    temp_dir = tempfile.mkdtemp()
    try:
        yield temp_dir
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

# LibreOffice bilan konvertatsiya funksiyalari (win32com o'rniga)
def convert_with_libreoffice(input_path: str, output_path: str, input_ext: str, output_ext: str = 'pdf') -> bool:
    """LibreOffice headless orqali fayl konvertatsiya qilish"""
    try:
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', output_ext,
            '--outdir', os.path.dirname(output_path),
            input_path
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=60)
        if result.returncode == 0:
            # LibreOffice chiqarilgan fayl nomini o'zgartiradi, to'g'rilaymiz
            expected_output = os.path.splitext(input_path)[0] + f'.{output_ext}'
            if os.path.exists(expected_output):
                os.rename(expected_output, output_path)
            return os.path.exists(output_path)
        else:
            logger.error(f"LibreOffice xato: {result.stderr.decode()}")
            return False
    except Exception as e:
        logger.error(f"LibreOffice konvertatsiyada xato: {e}")
        return False

# Boshqa funksiyalar (saqlangan, ozgina o'zgartirishlar bilan)
def format_file_size(size_bytes: int) -> str:
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.2f} TB"

def compress_image(input_path: str, output_path: str, quality: str) -> bool:
    try:
        quality_settings = {'high': 90, 'medium': 70, 'max': 50}
        with Image.open(input_path) as img:
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            img.save(output_path, "JPEG", quality=quality_settings[quality], optimize=True)
        return os.path.exists(output_path)
    except Exception as e:
        logger.error(f"Rasmni siqishda xato: {str(e)}")
        return False

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Start komandasi handleri"""
    load_channels_from_excel()
    
    user_id = update.effective_user.id
    user = update.effective_user
    first_name = user.first_name or "Foydalanuvchi"  # Use first_name, fallback to "Foydalanuvchi"
    
    users[user_id] = {
        'first_name': first_name,
        'username': f"@{user.username}" if user.username else 'N/A'
    }
    
    if user_id not in user_activity:
        user_activity[user_id] = []
    user_activity[user_id].append(datetime.now())
    
    if subscription_channels:
        is_subscribed = await check_all_subscriptions(user_id, context)
        if not is_subscribed:
            await send_subscription_request(update, context)
            return
    
    try:
        keyboard = [
            [InlineKeyboardButton("📄Word 🔄 PDF", callback_data='docx_pdf'),
             InlineKeyboardButton("📊Excel ➡ PDF", callback_data='excel_pdf')],
            [InlineKeyboardButton("🎤Slayd (PPTX) ➡ PDF", callback_data='ppt_pdf'),
             InlineKeyboardButton("🖼RASM(JPG) ➡ PDF", callback_data='jpg_pdf')],
            [InlineKeyboardButton("🔤Kiril 🔄 Lotin", callback_data='translate_file'),
             InlineKeyboardButton("🔲QR Kod Yasash", callback_data='qr_gen')],
            [InlineKeyboardButton("📷QR Kod Aniqlash", callback_data='qr_scan'),
             InlineKeyboardButton("💧PDF Suv belgi qo'yish", callback_data='pdf_watermark')],
            [InlineKeyboardButton("🔒PDF Parol qo'yish", callback_data='pdf_protect'),
             InlineKeyboardButton("🔓PDF Parolni olib tashlash", callback_data='pdf_unprotect')],
            [InlineKeyboardButton("📝Word ➡ Excel", callback_data='word_excel'),
             InlineKeyboardButton("🗂Arxivlovchi", callback_data='file_zipper')],
            [InlineKeyboardButton("📜Matn chiqarish", callback_data='ocr'),
             InlineKeyboardButton("🗜Fayl Siqish", callback_data='compress')],
            [InlineKeyboardButton("📄PDF Ajratish", callback_data='pdf_split'),
             InlineKeyboardButton("ℹYo'riqnoma", callback_data='about')]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
    
        welcome_text = f'''
    *Assalomu alaykum {first_name}!* 👋
Botni qayta ishga tushurish uchun /start buyrug'ini yuboring
O'zingizga kerakli bo'limni tanlang👇
        '''
    
        await update.message.reply_text(
            welcome_text,
            reply_markup=reply_markup,
            parse_mode="Markdown"
        )
    except Exception as e:
        logger.error(f"/start buyrug'ini qayta yuboring: {str(e)}")
        await update.message.reply_text("⚠️ Xato yuz berdi, iltimos /start ni qayta yuboring.")
        
async def delete_status_message(user_id: int, chat_id: int, context: ContextTypes.DEFAULT_TYPE) -> None:
    if user_id in user_data and user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
async def compress_file_temp(input_path: str, file_extension: str, quality: str) -> int:
    """Faylni vaqtincha siqib, hajmini hisoblash"""
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = os.path.join(temp_dir, f"temp_compressed{file_extension}")
            
            if file_extension == '.pdf':
                # PDF uchun
                pdf_doc = fitz.open(input_path)
                image_paths = []
                
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    zoom = {
                        'high': 300/72,
                        'medium': 150/72,
                        'max': 72/72
                    }[quality]
                    
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    
                    img_path = os.path.join(temp_dir, f"temp_page_{page_num}.jpg")
                    pix.save(img_path, "jpeg", quality={
                        'high': 90,
                        'medium': 70,
                        'max': 50
                    }[quality])
                    
                    image_paths.append(img_path)
                
                pdf_doc.close()
                
                # Rasmlardan PDF yaratish
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(image_paths))
                
                if not os.path.exists(output_path):
                    return 0
                
            else:
                # Rasm fayllari uchun
                with Image.open(input_path) as img:
                    if img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')
                    
                    quality_val = {
                        'high': 90,
                        'medium': 70,
                        'max': 50
                    }[quality]
                    
                    img.save(output_path, "JPEG", quality=quality_val, optimize=True)
            
            return os.path.getsize(output_path)
    
    except Exception as e:
        logger.error(f"Vaqtincha siqishda xato: {str(e)}")
        return 0
    

def parse_page_numbers(page_input: str, total_pages: int) -> List[int]:
    """Sahifa raqamlarini tahlil qiladi va ro'yxat sifatida qaytaradi"""
    try:
        pages = set()
        # Vergul bilan ajratilgan qismlarni tekshirish
        for part in page_input.replace(" ", "").split(","):
            part = part.strip()
            if not part:
                continue
                
            if "-" in part:
                # Diapazon (masalan, 4-10)
                start, end = map(int, part.split("-"))
                if start < 1 or end > total_pages or start > end:
                    raise ValueError(f"Noto'g'ri sahifa diapazoni: {part}")
                pages.update(range(start, end + 1))
            else:
                # Yagona sahifa (masalan, 3)
                page = int(part)
                if page < 1 or page > total_pages:
                    raise ValueError(f"Noto'g'ri sahifa raqami: {page}")
                pages.add(page)
                
        return sorted(list(pages))
    except ValueError as e:
        raise ValueError(f"Sahifa raqamlarini kiritishda xato: {str(e)}")
    except Exception as e:
        raise ValueError("Sahifa raqamlarini tahlil qilishda xato")

async def handle_pdf_split(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_split':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"split_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_split_pages'
        
        # PDF sahifalari sonini aniqlash
        pdf_doc = fitz.open(user_data[user_id].temp_file_path)
        total_pages = len(pdf_doc)
        pdf_doc.close()
        
        await msg.edit_text(
            f"✅ PDF fayl qabul qilindi! Jami sahifalar: {total_pages}\n\n"
            "Ajratmoqchi bo'lgan sahifalarni kiriting (masalan: 1,3,4 yoki 4-10 yoki 1,3,5,7-10):"
        )
        
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")
        if 'msg' in locals():
            await context.bot.delete_message(chat_id, msg.message_id)
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_split_pages(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_split_pages':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.text:
            await update.message.reply_text("❌ Iltimos, sahifa raqamlarini kiriting!")
            return
            
        page_input = update.message.text.strip()
        msg = await update.message.reply_text("⏳")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"split_output_{user_id}.pdf")
        
        # PDF faylni ochish
        pdf_doc = fitz.open(input_path)
        total_pages = len(pdf_doc)
        
        # Sahifa raqamlarini tahlil qilish
        try:
            page_numbers = parse_page_numbers(page_input, total_pages)
        except ValueError as e:
            await update.message.reply_text(f"❌ {str(e)}")
            return
        
        if not page_numbers:
            pdf_doc.close()
            raise ValueError("Hech qanday sahifa tanlanmadi❌")
        
        # Yangi PDF yaratish
        new_pdf = fitz.open()
        for page_num in page_numbers:
            new_pdf.insert_pdf(pdf_doc, from_page=page_num-1, to_page=page_num-1)
        
        new_pdf.save(output_path)
        new_pdf.close()
        pdf_doc.close()
        
        # Natijani yuborish
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                filename=f"split_{os.path.splitext(os.path.basename(input_path))[0]}.pdf",
                caption=f"✅ {len(page_numbers)} ta sahifa ajratildi: {page_input}\n🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
            user_data[user_id].active_module = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
        
# Excel to PDF handler (LibreOffice bilan)
async def handle_excel_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'excel_pdf':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
        
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, Excel fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.xlsx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.xls', '.xlsx']:
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, .xls yoki .xlsx fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring.")
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳")
        
        with temp_file_manager() as temp_dir:
            safe_file_name = re.sub(r'[^\w\-\.]', '_', file_name)
            input_path = os.path.join(temp_dir, safe_file_name)
            await file.download_to_drive(input_path)
            
            if not os.path.exists(input_path):
                raise FileNotFoundError(f"Fayl topilmadi: {input_path}")
            
            # LibreOffice bilan konvertatsiya
            output_path = os.path.join(temp_dir, f"{os.path.splitext(safe_file_name)[0]}.pdf")
            success = convert_with_libreoffice(input_path, output_path, file_extension)
            
            if not success or not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                raise Exception("PDF fayli yaratilmadi.")
            
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    caption=f"✅ {os.path.splitext(safe_file_name)[0]}.pdf\n\n🌐 @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        logger.error(f"Excel ni PDF ga o‘tkazishda xato: {str(e)}", exc_info=True)
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_ppt_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """PowerPoint faylni PDF ga aylantirish"""
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'ppt_pdf':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
        
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, PowerPoint fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "presentation.pptx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.ppt', '.pptx']:
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, .ppt yoki .pptx fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳ PPT fayl PDF ga aylantirilmoqda...")
        
        # Temp dir ishlatish (/tmp Render'da ishlaydi)
        with temp_file_manager() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
            
            # .ppt uchun avval .pptx ga konvertatsiya (LibreOffice)
            if file_extension == '.ppt':
                pptx_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pptx")
                if not convert_with_libreoffice(input_path, pptx_path, '.ppt', 'pptx'):
                    raise Exception("PPT ni PPTX ga konvertatsiya qilishda xato.")
                input_path = pptx_path  # Endi PPTX bilan ishlash
            
            # PDF ga konvertatsiya
            success = convert_pptx_to_pdf(input_path, output_path)
            
            if not success or not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                raise Exception("PDF fayli yaratilmadi yoki bo'sh.")
            
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    filename=f"{os.path.splitext(file_name)[0]}.pdf",
                    caption=f"✅ {os.path.splitext(file_name)[0]}.pdf\n\n🌐 @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        logger.error(f"PPT ni PDF ga o'tkazishda xato: {str(e)}", exc_info=True)
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")
        if 'msg' in locals():
            try:
                await context.bot.delete_message(chat_id, msg.message_id)
            except:
                pass
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    """PPTX ni PDF ga aylantirish: LibreOffice bilan (python-pptx PDF export qilmaydi)"""
    try:
        return convert_with_libreoffice(pptx_path, pdf_path, '.pptx')
    except Exception as e:
        logger.error(f"PPTX konvertatsiyada xato: {e}")
        return False

def convert_docx_to_pdf(docx_path: str, pdf_path: str) -> bool:
    """DOCX ni PDF ga aylantirish: docx2pdf bilan, fallback LibreOffice"""
    try:
        # Birinchi docx2pdf ni sinab ko'rish (tezroq)
        convert(docx_path, pdf_path)
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            return True
    except Exception as e:
        logger.warning(f"docx2pdf ishlamadi: {e}. LibreOffice ga o'tilmoqda.")
    
    # Fallback: LibreOffice
    return convert_with_libreoffice(docx_path, pdf_path, '.docx')

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Callback query handler – barcha tugmalar uchun"""
    query = update.callback_query
    await query.answer()
    
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    async with asyncio.Lock():
        if user_id not in user_data:
            user_data[user_id] = UserData()
    
    if query.data == 'excel_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'excel_pdf'
        message = await query.edit_message_text(
            "📤 *Excel faylingizni yuboring (.xls, .xlsx)*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Natija PDF formatida qaytariladi\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
            
    elif query.data == 'pdf_split':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_split'
            user_data[user_id].temp_file_path = None  # Qo'shimcha tozalash
        message = await query.edit_message_text(
            "📄 *PDF faylingizni yuboring*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Keyin ajratmoqchi bo'lgan sahifalarni kiriting (masalan: 1,3,4 yoki 4-10 yoki 1,3,5,7-10)\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
            
    elif query.data == 'ppt_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'ppt_pdf'
        message = await query.edit_message_text(
            "📤 *PowerPoint faylingizni yuboring (.ppt, .pptx)*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Natija PDF formatida qaytariladi\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'qr_gen':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'qr_gen'
        message = await query.edit_message_text(
            "🔲 *QR kod Yaratish*\n\n"
            "QR kodga aylantirmoqchi bo'lgan matn yoki URL manzilini yuboring:\n\n"
            "Misol: https://google.com yoki Convert Bot!\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'qr_scan':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'qr_scan'
        message = await query.edit_message_text(
            "📷 *QR kod Aniqlash*\n\n"
            "QR kodni o'qish uchun rasm yuboring:\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Rasm aniq va yorqin bo'lishi kerak\n"
            "- QR kod rasmdagi asosiy obyekt bo'lishi kerak\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_watermark':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_watermark'
        message = await query.edit_message_text(
            "💧 *PDF Suv belgisini qo'shish*\n\n"
            "Suv belgisi qo'shmoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga suv belgisi matnini so'raymiz\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_protect':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_protect'
        message = await query.edit_message_text(
            "🔒 *PDF faylga parol qo'yish*\n\n"
            "Parol qo'ymoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga parolni so'raymiz\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_unprotect':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_unprotect'
        message = await query.edit_message_text(
            "🔓 *PDF fayldan parolni olib tashlash*\n\n"
            "Parolini olib tashlamoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga parolni so'raymiz (agar ma'lum bo'lsa)\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'word_excel':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'word_excel'
        message = await query.edit_message_text(
            "📝 *Word dan Excelga o'tkazish*\n\n"
            "Excelga o'tkazmoqchi bo'lgan Word faylingizni yuboring:\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Faylda jadvallar bo'lishi kerak\n"
            "- Fayl hajmi 10MB dan oshmasin\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'compress':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'compress'
            user_data[user_id].compressed_sizes = {}  # Hajmlarni tozalash
            user_data[user_id].temp_file_path = None  # Fayl yo‘lini tozalash
        message = await query.edit_message_text(
            "📤 *JPG, PNG yoki PDF faylingizni yuboring*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Fayl siqilgandan so‘ng sifat tanlash imkoniyati beriladi\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug‘iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data in [COMPRESS_HIGH, COMPRESS_MEDIUM, COMPRESS_MAX]:
        compression_level = {
            COMPRESS_HIGH: 'high',
            COMPRESS_MEDIUM: 'medium',
            COMPRESS_MAX: 'max'
        }[query.data]
    
        await query.answer("Fayl siqilmoqda...")
    
        try:
            # Oldingi xabarni o'chirish
            await context.bot.delete_message(chat_id=chat_id, message_id=query.message.message_id)
        
            # Fayl ma'lumotlarini olish
            async with user_data[user_id].lock:
                if not user_data[user_id].temp_file_path or not os.path.exists(user_data[user_id].temp_file_path):
                    await context.bot.send_message(chat_id, "❌ Avval fayl yuboring yoki fayl topilmadi!")
                    return
            
                input_path = user_data[user_id].temp_file_path
                file_name = user_data[user_id].temp_file_name
                file_extension = os.path.splitext(file_name)[1].lower()
                temp_dir = user_data[user_id].temp_dir
        
            processing_msg = await context.bot.send_message(chat_id, "⏳")
        
            with temp_file_manager() as output_temp_dir:
                output_path = os.path.join(output_temp_dir, f"compressed_{file_name}")
            
                if file_extension == '.pdf':
                    success = compress_pdf(input_path, output_path, compression_level)  # compress_pdf funksiyasini oldingi kodingizdan nusxa ko'chiring
                else:
                    success = compress_image(input_path, output_path, compression_level)
                
                if not success:
                    raise Exception("Faylni siqish muvaffaqiyatsiz tugadi")
                    
                # Natijani yuborish
                with open(output_path, 'rb') as result_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=result_file,
                        filename=f"compressed_{file_name}",
                        caption=f"✅ Fayl siqildi ({compression_level})!\n🌐 @Convert_filesbot"
                    )
        
            await context.bot.delete_message(chat_id, processing_msg.message_id)
    
        except Exception as e:
            logger.error(f"Faylni siqishda xato: {str(e)}", exc_info=True)
            await context.bot.send_message(
                chat_id=chat_id,
                text=f"❌ Faylni siqishda xatolik yuz berdi: {str(e)}"
            )
    
        finally:
            # Tozalash
            async with user_data[user_id].lock:
                if 'input_path' in locals() and os.path.exists(input_path):
                    os.unlink(input_path)
                if 'temp_dir' in locals() and os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir, ignore_errors=True)
            
                user_data[user_id].temp_file_path = None
                user_data[user_id].temp_file_name = None
                user_data[user_id].compressed_sizes = {}
                user_data[user_id].temp_dir = None
        
            await asyncio.sleep(3)
            await return_to_main_menu(chat_id, context)

    # Qolgan button'lar (OCR, docx_pdf, jpg_pdf va boshqalar – oldingi kodingizdan nusxa)
    elif query.data == 'ocr':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'ocr'
        message = await query.edit_message_text(
            "📤 *JPG, PNG yoki PDF faylingizni yuboring*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Matn aniq va o‘qiladigan bo‘lishi kerak\n"
            "- Natija formatini keyin tanlaysiz\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data in [OCR_MESSAGE, OCR_TXT, OCR_DOCX]:
        async with user_data[user_id].lock:
            if not hasattr(user_data[user_id], 'ocr_text') or not user_data[user_id].ocr_text:
                await query.answer("Avval matn chiqarish uchun fayl yuboring!", show_alert=True)
                return
            output_format = query.data
            user_data[user_id].output_format = output_format
        
        try:
            await context.bot.delete_message(chat_id=chat_id, message_id=query.message.message_id)
        except Exception as e:
            logger.error(f"Format tanlash xabarini o‘chirishda xato: {e}")
        
        await query.answer("⏳")
        
        if output_format == OCR_MESSAGE:
            await context.bot.send_message(
                chat_id=chat_id,
                text=user_data[user_id].ocr_text[:4096],
                parse_mode="Markdown"
            )
        else:
            with temp_file_manager() as temp_dir:
                if output_format == OCR_TXT:
                    output_path = os.path.join(temp_dir, "extracted_text.txt")
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(user_data[user_id].ocr_text)
                    filename = "extracted_text.txt"
                else:
                    output_path = os.path.join(temp_dir, "extracted_text.docx")
                    doc = Document()
                    doc.add_paragraph(user_data[user_id].ocr_text)
                    doc.save(output_path)
                    filename = "extracted_text.docx"
                
                with open(output_path, 'rb') as result_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=result_file,
                        filename=filename,
                        caption=f"✅ Matn muvaffaqiyatli chiqarildi!\n🌐 @Convert_filesbot"
                    )
        
        async with user_data[user_id].lock:
            user_data[user_id].ocr_text = None
            user_data[user_id].output_format = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
    
    elif query.data == 'docx_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'docx_pdf'
        message = await query.edit_message_text(
            "📤 *PDF yoki DOCX faylingizni yuboring*\n\n"
            "❗ Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Faqat .docx yoki .pdf formatdagi fayllar qabul qilinadi\n"
            "- PDF fayllar DOCX ga, DOCX fayllar PDF ga aylantiriladi\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'jpg_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'jpg_pdf'
            user_data[user_id].images = []
            user_data[user_id].waiting_for_images = True
        message = await query.edit_message_text(
            "📤 PDF qilish kerak bo'lgan rasmlarni barchasini yuboring. "
            "Barcha rasmlarni yuborgach,\n"
            "*PDF qilish* tugmasini bosing.\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'translate_file':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'kiril_lotin'
        message = await query.edit_message_text(
            "Iltimos, tarjima qilish uchun fayl yuboring:\n\n"
            "📝 Qo'llab-quvvatlanadigan formatlar:\n"
            "- Word hujjatlari (.doc, .docx)\n"
            "- Excel jadvallari (.xls, .xlsx)\n"
            "- PowerPoint prezentatsiyalari (.ppt, .pptx)\n"
            "- PDF fayllari (.pdf)\n\n"
            "Bot avtomatik tarjima tilini aniqlaydi!\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'file_zipper':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'file_zipper'
            user_data[user_id].files = []
        
        keyboard = [
            [InlineKeyboardButton("🗂 Arxivlash (0)", callback_data=ZIP_BUTTON)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        message = await query.edit_message_text(
            "📁 <b>Zip Botga xush kelibsiz!</b>\n\n"
            "🔹 Istalgan formatdagi fayllarni yuboring\n"
            "🔹 Barcha fayllar ro'yxatga qo'shiladi\n"
            "🔹 Pastdagi tugma orqali zip faylni oling\n\n"
            "📂 Joriy fayllar ro'yxati:\n"
            "Hozircha fayllar yo'q\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            reply_markup=reply_markup,
            parse_mode=ParseMode.HTML
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'about':
        await query.edit_message_text(
            """
     <b>YO'RIQNOMA!</b> 👋
    Quyidagi funksiyalar mavjud:
- 📄 <b>Word 🔄 PDF</b> - Word faylini PDF ga yoki teskarisiga o'tkazish
- 📊 <b>Excel ➡ PDF</b> - Excel jadvallarini PDFga aylantirish
- 🎤 <b>Slayd(PPTX) ➡ PDF</b> - PowerPoint fayllarini PDFga aylantirish
- 🖼 <b>RASM(JPG)➡ PDF</b> - Yuborilgan Rasmlarni PDF qilish
- 🔤 <b>Kiril 🔄 Lotin</b> - Fayllarni avtomatik Lotin Kirilga o'tkazish
- 🔲 <b>QR Kod Yasash</b> - Matn yoki URLdan QR kod yaratish
- 📷 <b>QR Kod Aniqlash</b> - Rasm orqali QR koddan ma'lumot olish
- 💧 <b>PDF Suv belgi qo'yyish</b> - PDF fayllarga suv belgisi qo'shish
- 🔒 <b>PDF Parol qo'yish</b> - PDF fayllarga himoya qo'yish
- 🔓 <b>PDF Parolni olib tashlash</b> - Parolli PDFlarni ochish
- 📝 <b>Word ➡ Excel</b> - Word hujjatidagi jadvallarni Excelga o'tkazish
- 📜 <b>Matn chiqarish </b> - Siz yuborgan rasmdan matnni chiqarish
- 🗜 <b>Fayl Siqish</b> - Rasm yoki PDF formatdagi faylni siqish
- 🗂 <b>Arxivlovchi</b> - Fayllarni arxivlash
- 📄 <b>PDF Ajratish</b> - PDF faylni sahifalar bo'yicha ajratish

<b>Aloqa:</b> @Dilxush_Bahodirov
<b>ℹ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!</b>
            """,
            parse_mode=ParseMode.HTML
        )
    
    elif query.data == ZIP_BUTTON:
        await zip_files(update, context)  # zip_files funksiyasini oldingi kodingizdan nusxa ko'chiring
    
    elif query.data == CLEAR_BUTTON:
        await clear_files(update, context)  # clear_files funksiyasini oldingi kodingizdan nusxa ko'chiring
    
    elif query.data == 'create_pdf':
        await create_pdf(update, context)  # create_pdf funksiyasini oldingi kodingizdan nusxa ko'chiring


async def handle_qr_gen(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'qr_gen':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.text:
            await update.message.reply_text("❌ Iltimos, QR kodga aylantirmoqchi bo'lgan matn yuboring!")
            return
            
        text = update.message.text
        msg = await update.message.reply_text("⏳")
        
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=4,
        )
        qr.add_data(text)
        qr.make(fit=True)
        
        img = qr.make_image(fill_color="black", back_color="white")
        
        with BytesIO() as bio:
            img.save(bio, 'PNG')
            bio.seek(0)
            await update.message.reply_photo(
                photo=bio,
                caption=f"✅ QR kod yaratildi!\nMatn: {text[:50]}{'...' if len(text) > 50 else ''}\n🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ QR kod yaratishda xatolik: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_qr_scan(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'qr_scan':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.photo:
            await update.message.reply_text("❌ Iltimos, QR kod rasmini yuboring!")
            return
            
        photo = update.message.photo[-1]
        msg = await update.message.reply_text("⏳")
        
        file = await photo.get_file()
        with BytesIO() as bio:
            await file.download_to_memory(out=bio)
            bio.seek(0)
            
            img = Image.open(bio)
            img_array = np.array(img)
            
            detector = cv2.QRCodeDetector()
            data, _, _ = detector.detectAndDecode(img_array)
            
            if data:
                await update.message.reply_text(
                    f"✅ QR kod mazmuni:\n\n{data}\n\n🌐 @Convert_filesbot",
                    disable_web_page_preview=True
                )
            else:
                await update.message.reply_text("❌ QR kod aniqlanmadi yoki o'qib bo'lmadi!")
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ QR kodni o'qishda xatolik: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_watermark(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_watermark':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳ PDF fayl yuklanmoqda...")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"watermark_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_watermark_text'
        
        await msg.edit_text("✅ PDF fayl qabul qilindi!\n\nEndi suv belgisi uchun matn yuboring:")
        
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")

async def handle_pdf_watermark_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_watermark_text':
        return
        
    try:
        if not update.message.text:
            await update.message.reply_text("❌ Iltimos, suv belgisi matnini yuboring!")
            return
            
        watermark_text = update.message.text
        msg = await update.message.reply_text("⏳")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"watermarked_{user_id}.pdf")
        
        # Watermark qo'shish
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFont("Helvetica", 50)
        can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
        can.rotate(45)
        can.drawString(100, 100, watermark_text)
        can.save()
        
        packet.seek(0)
        watermark_pdf = fitz.open("pdf", packet.getvalue())
        
        input_pdf = fitz.open(input_path)
        
        for page in input_pdf:
            page.show_pdf_page(page.rect, watermark_pdf, 0)
        
        input_pdf.save(output_path)
        input_pdf.close()
        watermark_pdf.close()
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"✅ Suv belgisi qo'shilgan PDF\nMatn: {watermark_text}\n🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ Suv belgisi qo'shishda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_protect(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_protect':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"protect_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_protect_password'
        
        await msg.edit_text("✅ PDF fayl qabul qilindi!\n\nEndi parolni yuboring (6-32 belgi):")
        
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")

async def handle_pdf_protect_password(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_protect_password':
        return
        
    try:
        if not update.message.text:
            await update.message.reply_text("❌ Iltimos, parol yuboring!")
            return
            
        password = update.message.text
        if len(password) < 6 or len(password) > 32:
            await update.message.reply_text("❌ Parol 6-32 belgidan iborat bo'lishi kerak!")
            return
            
        msg = await update.message.reply_text("⏳")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"protected_{user_id}.pdf")
        
        pdf_reader = PyPDF2.PdfReader(input_path)
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        pdf_writer.encrypt(password)
        
        with open(output_path, "wb") as f:
            pdf_writer.write(f)
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"✅ Parol bilan himoyalangan PDF\n🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ Parol qo'yishda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_unprotect(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_unprotect':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, parolli PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"unprotect_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_unprotect_password'
        
        await msg.edit_text("✅ PDF fayl qabul qilindi!\n\nAgar parol ma'lum bo'lsa, yuboring (agar parolni bilmasangiz, 'keyin' yozing):")
        
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")

async def handle_pdf_unprotect_password(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_unprotect_password':
        return
        
    try:
        password = update.message.text if update.message.text.lower() != 'keyin' else None
        
        msg = await update.message.reply_text("⏳")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"unprotected_{user_id}.pdf")
        
        pdf_reader = PyPDF2.PdfReader(input_path)
        if pdf_reader.is_encrypted:
            if password:
                pdf_reader.decrypt(password)
            else:
                # Parolni bilmasa, oddiy usul bilan ochishga harakat qilish
                try:
                    pdf_reader.decrypt("")
                except:
                    raise Exception("Parol talab qilinadi va siz parolni yubormagansiz")
        
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        with open(output_path, "wb") as f:
            pdf_writer.write(f)
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"✅ Parolsiz PDF\n🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"❌ Parolni olib tashlashda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_word_excel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'word_excel':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")    
    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, Word fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.docx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.doc', '.docx']:
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, Word fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.xlsx")
            
            # Convert .doc to .docx if necessary
            if file_extension == '.doc':
                temp_docx_path = os.path.join(temp_dir, "temp.docx")
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                doc = word.Documents.Open(input_path)
                doc.SaveAs(temp_docx_path, FileFormat=16)  # Save as DOCX
                doc.Close()
                word.Quit()
                pythoncom.CoUninitialize()
                input_path = temp_docx_path
            
            # Read the Word document
            doc = Document(input_path)
            
            # Create a new Excel workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Asosiy"
            
            # Track current row in Excel
            current_row = 1
            
            # Extract all content including tables
            for element in doc.element.body:
                if element.tag.endswith('tbl'):  # Jadval topildi
                    table = docx.table.Table(element, doc)
                    for row in table.rows:
                        for col_idx, cell in enumerate(row.cells):
                            # Jadval kataklarini Excelga yozish
                            ws.cell(row=current_row, column=col_idx+1).value = cell.text
                        current_row += 1
                    current_row += 1  # Jadval orasida bo'sh qator
                
                elif element.tag.endswith('p'):  # Oddiy paragraf
                    paragraph = docx.text.paragraph.Paragraph(element, doc)
                    if paragraph.text.strip():  # Faqat matn bor paragraflarni olish
                        ws.cell(row=current_row, column=1).value = paragraph.text
                        current_row += 1
            
            # Adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(cell.value)
                    except:
                        pass
                adjusted_width = (max_length + 2) * 1.2
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Save the Excel file
            wb.save(output_path)
            
            # Send the result
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    filename=f"{os.path.splitext(file_name)[0]}.xlsx",
                    caption=f"✅ {os.path.splitext(file_name)[0]}.xlsx\n\n🌐 @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await msg.edit_text(f"❌ Xatolik yuz berdi: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_compress(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    async with user_data[user_id].lock:
        if user_id not in user_data or user_data[user_id].active_module != 'compress':
            return
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    try:
        # Fayl yuklash va tekshirish
        if update.message.document:
            document = update.message.document
            file_name = document.file_name or "document.pdf"
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension != '.pdf':
                await update.message.reply_text("❌ Faqat PDF fayllar qabul qilinadi!")
                return
            file = await document.get_file()
        elif update.message.photo:
            file = await update.message.photo[-1].get_file()
            file_name = "photo.jpg"
            file_extension = '.jpg'
        else:
            await update.message.reply_text("❌ Iltimos, JPG yoki PDF fayl yuboring!")
            return
        
        if file.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        msg = await update.message.reply_text("⏳")
        
        # Faylni yuklash
        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, file_name)
        await file.download_to_drive(input_path)
        
        # PDF uchun maxsus tekshiruv
        if file_extension == '.pdf':
            try:
                with fitz.open(input_path) as doc:
                    if doc.is_encrypted:
                        await msg.edit_text("❌ PDF fayl parol bilan himoyalangan. Iltimos, parolsiz PDF yuboring.")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return
            except Exception as e:
                await msg.edit_text(f"❌ PDF faylni ochib bo'lmadi: {str(e)}")
                shutil.rmtree(temp_dir, ignore_errors=True)
                return
        
        # Fayl hajmlarini hisoblash
        sizes = {}
        for quality in ['high', 'medium', 'max']:
            with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as temp_file:
                output_path = temp_file.name
            
            if file_extension == '.pdf':
                success = compress_pdf(input_path, output_path, quality)
            else:
                success = compress_image(input_path, output_path, quality)
            
            if success:
                sizes[quality] = os.path.getsize(output_path)
                os.unlink(output_path)
            else:
                sizes[quality] = 0
        
        # Agar siqish muvaffaqiyatsiz bo'lsa
        if all(size == 0 for size in sizes.values()):
            raise Exception("Faylni siqish jarayonida xatolik yuz berdi")
        
        # Ma'lumotlarni saqlash
        async with user_data[user_id].lock:
            user_data[user_id].compressed_sizes = sizes
            user_data[user_id].temp_file_path = input_path
            user_data[user_id].temp_file_name = file_name
            user_data[user_id].temp_dir = temp_dir
        
        # Inline tugmalar
        keyboard = [
            [InlineKeyboardButton(f"Yuqori ({format_file_size(sizes['high'])})", callback_data=COMPRESS_HIGH)],
            [InlineKeyboardButton(f"Oʻrta ({format_file_size(sizes['medium'])})", callback_data=COMPRESS_MEDIUM)],
            [InlineKeyboardButton(f"Maksimal ({format_file_size(sizes['max'])})", callback_data=COMPRESS_MAX)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await msg.edit_text(
            "✅ Fayl tahlil qilindi!\n"
            "Quyidagi siqish darajalaridan birini tanlang:\n"
            f"- Yuqori: {format_file_size(sizes['high'])}\n"
            f"- Oʻrta: {format_file_size(sizes['medium'])}\n"
            f"- Maksimal: {format_file_size(sizes['max'])}\n\n"
            "ℹ Dasturni qayta ishga tushurish uchun /start buyrugʻiga bosing!",
            reply_markup=reply_markup,
            parse_mode="Markdown"
        )
        
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = msg.message_id
    
    except Exception as e:
        logger.error(f"Faylni siqishda xato: {str(e)}", exc_info=True)
        error_msg = f"❌ Xatolik yuz berdi: {str(e)}"
        if "PDF" in str(e):
            error_msg += "\n\n⚠️ Iltimos, quyidagilarni tekshiring:"
            error_msg += "\n1. PDF fayl buzilmaganligiga ishonch hosil qiling"
            error_msg += "\n2. Fayl parol bilan himoyalanmaganligiga ishonch hosil qiling"
            error_msg += "\n3. Fayl hajmi 10MB dan oshmasligiga ahamiyat bering"
        
        if 'msg' in locals():
            await msg.edit_text(error_msg)
        else:
            await update.message.reply_text(error_msg)
        
        # Tozalash
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
        async with user_data[user_id].lock:
            user_data[user_id].compressed_sizes = {}
            user_data[user_id].temp_file_path = None
            user_data[user_id].temp_file_name = None
            user_data[user_id].temp_dir = None

def compress_pdf(input_path: str, output_path: str, quality: str) -> bool:
    """PDF faylni siqib, yangi faylga saqlaydi"""
    try:
        # Siqish parametrlari
        quality_settings = {
            'high': {'dpi': 300, 'jpeg_quality': 90},
            'medium': {'dpi': 150, 'jpeg_quality': 70},
            'max': {'dpi': 72, 'jpeg_quality': 50}
        }
        settings = quality_settings[quality]

        # PDF faylni ochish
        pdf = fitz.open(input_path)
        image_paths = []

        # Har bir sahifani rasmga aylantirish
        for page in pdf:
            pix = page.get_pixmap(matrix=fitz.Matrix(settings['dpi']/72, settings['dpi']/72))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Rasmni siqish
            img_buffer = BytesIO()
            img.save(img_buffer, format="JPEG", quality=settings['jpeg_quality'], optimize=True)
            
            # Vaqtinchalik faylga saqlash
            temp_img_path = os.path.join(tempfile.gettempdir(), f"temp_page_{len(image_paths)}.jpg")
            with open(temp_img_path, "wb") as f:
                f.write(img_buffer.getvalue())
            image_paths.append(temp_img_path)

        pdf.close()

        # Rasmlardan yangi PDF yaratish
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))

        # Vaqtinchalik fayllarni tozalash
        for img_path in image_paths:
            if os.path.exists(img_path):
                os.unlink(img_path)

        return os.path.exists(output_path)

    except Exception as e:
        logger.error(f"PDF siqishda xato: {str(e)}")
        # Vaqtinchalik fayllarni tozalash
        if 'image_paths' in locals():
            for img_path in image_paths:
                if os.path.exists(img_path):
                    os.unlink(img_path)
        return False
    
async def handle_ocr(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    async with user_data[user_id].lock:
        if user_id not in user_data or user_data[user_id].active_module != 'ocr':
            return
    
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    try:
        if update.message.document:
            document = update.message.document
            file_name = document.file_name or "document"
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension not in ['.pdf']:
                await update.message.reply_text("❌ Faqat PDF fayllar qabul qilinadi!")
                return
            file = await document.get_file()
        elif update.message.photo:
            file = await update.message.photo[-1].get_file()
            file_name = "photo.jpg"
            file_extension = '.jpg'
        else:
            await update.message.reply_text("❌ Iltimos, JPG, PNG yoki PDF fayl yuboring!")
            return
        
        if file.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        msg = await update.message.reply_text("⏳️")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            extracted_text = ""
            if file_extension == '.pdf':
                pdf_doc = fitz.open(input_path)
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_path = os.path.join(temp_dir, f"page_{page_num}.png")
                    pix.save(img_path)
                    extracted_text += pytesseract.image_to_string(Image.open(img_path), lang='uzb+eng')
                pdf_doc.close()
            else:
                extracted_text = pytesseract.image_to_string(Image.open(input_path), lang='uzb+eng')
            
            if not extracted_text.strip():
                await msg.edit_text("❌ Faylda matn topilmadi yoki o‘qib bo‘lmadi!")
                return
            
            async with user_data[user_id].lock:
                user_data[user_id].ocr_text = extracted_text
            
            keyboard = [
                [InlineKeyboardButton("Botga xabar sifatida", callback_data=OCR_MESSAGE)],
                [InlineKeyboardButton("TEXT fayl sifatida", callback_data=OCR_TXT)],
                [InlineKeyboardButton("Word fayl sifatida", callback_data=OCR_DOCX)]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await msg.edit_text(
                "✅ Matn muvaffaqiyatli chiqarildi!\n"
                "Qanday formatda olishni xohlaysiz?",
                reply_markup=reply_markup
            )
    
    except Exception as e:
        await msg.edit_text(f"❌ Xatolik yuz berdi: {str(e)}")

def convert_docx_to_pdf_alternative(docx_path, pdf_path):
    try:
        if sys.platform == "win32":
            libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        else:
            libreoffice_path = "soffice"
        
        command = [
            libreoffice_path,
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            os.path.dirname(pdf_path),
            docx_path
        ]
        
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return True
    
    except (subprocess.CalledProcessError, FileNotFoundError):
        try:
            doc = Document(docx_path)
            pdf_doc = fitz.open()
            
            for para in doc.paragraphs:
                if para.text.strip():
                    page = pdf_doc.new_page()
                    page.insert_text((50, 50), para.text)
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            return True
            
        except Exception as e:
            logger.error(f"Alternative DOCX to PDF conversion failed: {str(e)}")
            return False
        
async def handle_docx_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle DOC/DOCX to PDF and PDF to DOCX conversion"""
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id

    if user_id not in user_data or user_data[user_id].active_module != 'docx_pdf':
        return

    # Delete previous status message if it exists
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")

    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, DOC, DOCX yoki PDF fayl yuboring!")
            return

        document = update.message.document
        file_name = document.file_name or "document"
        file_extension = os.path.splitext(file_name)[1].lower()

        if file_extension not in ['.doc', '.docx', '.pdf']:
            await update.message.reply_text("❌ Noto'g'ri format! Iltimos, .doc, .docx yoki .pdf fayl yuboring.")
            return

        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"❌ Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return

        file = await update.message.document.get_file()
        msg = await update.message.reply_text("⏳ Fayl qayta ishlanmoqda...")
        user_data[user_id].status_message_id = msg.message_id

        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)

            if file_extension in ['.doc', '.docx']:
                # Convert DOC/DOCX to PDF
                output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")

                try:
                    pythoncom.CoInitialize()
                    if file_extension == '.docx':
                        convert(input_path, output_path)
                    else:
                        word = None
                        doc = None
                        try:
                            word = win32com.client.Dispatch("Word.Application")
                            word.Visible = False
                            word.DisplayAlerts = False
                            doc = word.Documents.Open(input_path)
                            doc.SaveAs(output_path, FileFormat=17)  # 17 is PDF format
                        finally:
                            if doc:
                                doc.Close(SaveChanges=False)
                            if word:
                                word.Quit()
                except Exception as com_error:
                    logger.error(f"COM xatosi: {str(com_error)}")
                    raise Exception(f"COM xatosi: {str(com_error)}")
                finally:
                    pythoncom.CoUninitialize()

                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    raise Exception("PDF fayli yaratilmadi yoki bo'sh.")

                # Delete processing message before sending result
                await context.bot.delete_message(chat_id, msg.message_id)
                user_data[user_id].status_message_id = None

                with open(output_path, 'rb') as result_file:
                    await update.message.reply_document(
                        document=result_file,
                        filename=f"{os.path.splitext(file_name)[0]}.pdf",
                        caption=f"✅ Fayl PDF ga aylantirildi!\n🌐 @Convert_filesbot"
                    )

            elif file_extension == '.pdf':
                # Convert PDF to DOCX
                output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.docx")
                
                converter = Converter(input_path)
                converter.convert(output_path)
                converter.close()

                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    raise Exception("DOCX fayli yaratilmadi yoki bo'sh.")

                # Delete processing message before sending result
                await context.bot.delete_message(chat_id, msg.message_id)
                user_data[user_id].status_message_id = None

                with open(output_path, 'rb') as result_file:
                    await update.message.reply_document(
                        document=result_file,
                        filename=f"{os.path.splitext(file_name)[0]}.docx",
                        caption=f"✅ Fayl DOCX ga aylantirildi!\n🌐 @Convert_filesbot"
                    )

    except Exception as e:
        logger.error(f"Fayl konvertatsiyasida xato: {str(e)}")
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
                user_data[user_id].status_message_id = None
            except Exception as delete_error:
                logger.error(f"Xabarni o'chirishda xato: {delete_error}")
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")

    finally:
        async with user_data[user_id].lock:
            user_data[user_id].active_module = None
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'jpg_pdf':
        return
    
    if not user_data[user_id].waiting_for_images:
        return
    
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    photo = update.message.photo[-1]
    
    try:
        photo_file = await photo.get_file()
        photo_bytes = BytesIO()
        await photo_file.download_to_memory(out=photo_bytes)
        photo_bytes.seek(0)
        
        with Image.open(photo_bytes) as img:
            img.verify()
        
        photo_bytes.seek(0)
        user_data[user_id].images.append(photo_bytes)
    except Exception as e:
        logger.error(f"Rasmni qayta ishlashda xato: {e}")
        await update.message.reply_text(f"Rasmni qayta ishlashda xato: {e}")
        return
    
    keyboard = [[InlineKeyboardButton("PDF qilish", callback_data='create_pdf')]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    message = await update.message.reply_text(
        f"🖼 Jami: {len(user_data[user_id].images)} ta rasm qabul qilindi.\n"
        "Yana rasm yuborishingiz mumkin yoki PDF yaratish tugmasini bosing🆗.",
        reply_markup=reply_markup
    )
    user_data[user_id].status_message_id = message.message_id

async def create_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    if user_id not in user_data or not user_data[user_id].images:
        await query.answer("Siz hech qanday rasm yubormagansiz!", show_alert=True)
        return
    
    await query.answer("⏳️")
    
    try:
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        temp_files = []
        
        for img_bytes in user_data[user_id].images:
            try:
                img_bytes.seek(0)
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                    temp_path = temp_img.name
                    temp_files.append(temp_path)
                    
                    with Image.open(img_bytes) as img:
                        if img.mode in ('RGBA', 'LA', 'P'):
                            img = img.convert('RGB')
                        img.save(temp_path, "JPEG", quality=90)
                    
                    pdf.add_page()
                    pdf.image(temp_path, x=10, y=10, w=190)
            except Exception as e:
                logger.error(f"Rasmni PDFga qo'shishda xato: {e}")
                continue
        
        if pdf.pages:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                pdf_path = temp_pdf.name
                pdf.output(pdf_path)
                
                with open(pdf_path, 'rb') as pdf_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=pdf_file,
                        filename='images.pdf',
                        caption=f"Sizning {len(user_data[user_id].images)} ta rasmingizdan PDF fayl yaratildi✅ \n🌐 @Convert_filesbot"
                    )
            
            os.unlink(pdf_path)
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
    except Exception as e:
        logger.error(f"PDF yaratishda xato: {e}")
        await context.bot.send_message(
            chat_id=chat_id,
            text=f"PDF yaratishda xatolik yuz berdi: {str(e)}"
        )
    finally:
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except Exception as e:
                logger.error(f"Vaqtinchalik faylni o'chirishda xato: {e}")
        
        if user_id in user_data:
            user_data[user_id].images = []
            user_data[user_id].waiting_for_images = False

async def handle_translation_doc(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'kiril_lotin':
        return

    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            user_data[user_id].status_message_id = None
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    progress_msg = None
    temp_dir = None
    output_path = None
    file_path = None

    try:
        if not update.message.document:
            await update.message.reply_text("❌ Iltimos, tarjima qilish uchun fayl yuboring!")
            return
        
        document = update.message.document
        file_name = document.file_name or "document"
        file_extension = os.path.splitext(file_name)[1].lower()

        if file_extension not in ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf']:
            await update.message.reply_text(
                "⚠️ Iltimos, quyidagi formatlardagi fayllarni yuboring:\n"
                "📝 Word (.doc, .docx)\n"
                "📊 Excel (.xls, .xlsx)\n"
                "🎤 PowerPoint (.ppt, .pptx)\n"
                "📄 PDF (.pdf)"
            )
            return

        progress_msg = await update.message.reply_text("⏳ Fayl yuklanmoqda...")

        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, file_name)
        
        file = await context.bot.get_file(document.file_id)
        await file.download_to_drive(file_path)
        
        original_extension = file_extension
        temp_converted_path = None
        
        if file_extension == '.pdf':
            await progress_msg.edit_text("⏳ PDF fayl DOCX ga o'tkazilmoqda...")
            temp_docx_path = os.path.join(temp_dir, f"temp_{os.path.splitext(file_name)[0]}.docx")
            
            try:
                cv = Converter(file_path)
                cv.convert(temp_docx_path, start=0, end=None, keep_layout=True, recognize=True)
                cv.close()
                
                if not os.path.exists(temp_docx_path) or os.path.getsize(temp_docx_path) < 1024:
                    await progress_msg.edit_text("ℹ️ Alternativ usul bilan qayta ishlanmoqda...")
                    pdf_doc = fitz.open(file_path)
                    doc = docx.Document()
                    
                    first_page = pdf_doc.load_page(0)
                    section = doc.sections[0]
                    section.page_width = Inches(first_page.rect.width / 72)
                    section.page_height = Inches(first_page.rect.height / 72)
                    
                    for page_num in range(len(pdf_doc)):
                        page = pdf_doc.load_page(page_num)
                        blocks = page.get_text("dict")["blocks"]
                        
                        for block in blocks:
                            if "lines" in block:
                                for line in block["lines"]:
                                    for span in line["spans"]:
                                        p = doc.add_paragraph()
                                        run = p.add_run(span["text"])
                                        font = run.font
                                        font.name = span["font"]
                                        font.size = Pt(span["size"])
                            
                            elif "image" in block:
                                try:
                                    img_info = block["image"]
                                    image_stream = BytesIO(img_info["image"])
                                    doc.add_picture(image_stream, 
                                                  width=Inches(img_info["width"] / 72),
                                                  height=Inches(img_info["height"] / 72))
                                except Exception:
                                    pass
                        
                        if page_num < len(pdf_doc) - 1:
                            doc.add_page_break()
                    
                    doc.save(temp_docx_path)
                    pdf_doc.close()
                
                file_path = temp_docx_path
                original_extension = '.docx'
                
            except Exception as e:
                logger.error(f"PDFni DOCXga o'tkazishda xato: {e}")
                raise Exception(f"PDF faylni DOCX formatiga o'tkazishda xatolik: {str(e)}")

        if file_extension in ['.doc', '.xls', '.ppt']:
            convert_functions = {
                '.doc': (convert_doc_to_docx, '.docx'),
                '.xls': (convert_xls_to_xlsx, '.xlsx'),
                '.ppt': (convert_ppt_to_pptx, '.pptx')
            }
            
            convert_func, new_ext = convert_functions[file_extension]
            await progress_msg.edit_text(f"⏳ {file_extension.upper()} fayl {new_ext.upper()} ga o'tkazilmoqda...")
            
            temp_converted_path = os.path.join(temp_dir, f"converted{new_ext}")
            if not convert_func(file_path, temp_converted_path):
                raise Exception(f"{file_extension.upper()} -> {new_ext.upper()} konvertatsiyada xatolik")
            
            file_path = temp_converted_path
            original_extension = new_ext

        await progress_msg.edit_text("⏳ Fayl tahlil qilinmoqda...")
        sample_text = extract_sample_text(file_path, original_extension)
        
        if not sample_text.strip():
            if original_extension in ['.xls', '.xlsx']:
                sample_text = "1"
            else:
                raise Exception("Faylda tarjima qilinadigan matn topilmadi")

        script = detect_script(sample_text)
        if script == 'cyrillic':
            translation_dict = CYRILLIC_TO_LATIN
            direction = 'cyr_lat'
            detected_script = "Kirill"
            target_script = "Lotin"
        else:
            translation_dict = LATIN_TO_CYRILLIC
            direction = 'lat_cyr'
            detected_script = "Lotin"
            target_script = "Kirill"

        output_filename = f"tarjima_{os.path.splitext(file_name)[0]}{original_extension}"
        output_path = os.path.join(temp_dir, output_filename)
        
        await progress_msg.edit_text(f"⏳ {detected_script} ➡️ {target_script} tarjima qilinmoqda...")
        
        if not translate_file(file_path, output_path, translation_dict, direction, original_extension):
            raise Exception("Tarjima jarayonida xatolik yuz berdi")

        if file_extension == '.pdf':
            pdf_output_path = os.path.join(temp_dir, f"tarjima_{os.path.splitext(file_name)[0]}.pdf")
            await progress_msg.edit_text("⏳ Tarjima qilingan DOCX PDF ga o'tkazilmoqda...")
            
            try:
                convert(output_path, pdf_output_path)
                output_path = pdf_output_path
                output_filename = f"tarjima_{os.path.splitext(file_name)[0]}.pdf"
            except Exception as e:
                logger.error(f"DOCXni PDFga o'tkazishda xato: {e}")
                raise Exception(f"Tarjima qilingan DOCXni PDFga o'tkazishda xatolik: {str(e)}")

        await progress_msg.edit_text("✅ Tarjima tugallandi! Fayl yuborilmoqda...")
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                filename=output_filename,
                caption=f"✅ {detected_script} ➡️ {target_script} tarjima qilindi\n"
                        f"📄 Original fayl: {file_name}\n"
                        f"🌐 @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, progress_msg.message_id)
        
    except Exception as e:
        logger.error(f"Faylni qayta ishlashda xatolik: {e}", exc_info=True)
        
        error_msg = "❌ Faylni qayta ishlashda xatolik yuz berdi.\n"
        
        if "The file is corrupt" in str(e):
            error_msg += "Fayl buzilgan yoki parol bilan himoyalangan.\n"
        elif "The password is incorrect" in str(e):
            error_msg += "Fayl parol bilan himoyalangan.\n"
        elif "Faylda tarjima qilinadigan matn topilmadi" in str(e):
            error_msg += "Faylda tarjima qilinadigan matn topilmadi.\n"
        else:
            error_msg += f"Xato tafsilotlari: {str(e)}\n"
        
        error_msg += "\nIltimos, quyidagilarni tekshiring:\n"
        error_msg += "1. Fayl formati qo'llab-quvvatlanadimi?\n"
        error_msg += "2. Fayl buzilmaganligiga ishonch hosil qiling\n"
        error_msg += "3. Fayl paroli bilan himoyalanmaganligiga ishonch hosil qiling\n\n"
        error_msg += "Agar muammo davom etsa, /start buyrug'i orqali qayta urinib ko'ring."
        
        if progress_msg:
            try:
                await progress_msg.edit_text(error_msg)
            except:
                await update.message.reply_text(error_msg)
        else:
            await update.message.reply_text(error_msg)
        
    finally:
        try:
            if output_path and os.path.exists(output_path):
                os.unlink(output_path)
            if file_path and os.path.exists(file_path):
                os.unlink(file_path)
            if temp_dir and os.path.exists(temp_dir):
                for filename in os.listdir(temp_dir):
                    file_path = os.path.join(temp_dir, filename)
                    try:
                        if os.path.isfile(file_path):
                            os.unlink(file_path)
                    except Exception as e:
                        logger.error(f"{file_path} ni o'chirishda xato: {e}")
                os.rmdir(temp_dir)
        except Exception as e:
            logger.error(f"Tozalashda xato: {e}")

    await asyncio.sleep(3)
    await return_to_main_menu(chat_id, context)
            
async def handle_zip_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'file_zipper':
        return
    
    try:
        if update.message.document:
            file = await update.message.document.get_file()
            original_name = update.message.document.file_name or f"file_{len(user_data[user_id].files)}"
        elif update.message.photo:
            photo = update.message.photo[-1]
            file = await photo.get_file()
            original_name = f"🖼 Rasm_{len(user_data[user_id].files)}.jpg"
        elif update.message.video:
            file = await update.message.video.get_file()
            original_name = update.message.video.file_name or f"🎬 Video_{len(user_data[user_id].files)}.mp4"
        elif update.message.audio:
            file = await update.message.audio.get_file()
            original_name = update.message.audio.file_name or f"🎵 Audio_{len(user_data[user_id].files)}.mp3"
        else:
            await update.message.reply_text("❌ Qo'llab-quvvatlanmaydigan fayl turi!")
            return
        
        file_path = f"temp_{user_id}_{len(user_data[user_id].files)}"
        await file.download_to_drive(file_path)
        
        unique_name = generate_unique_name(original_name)
        user_data[user_id].files.append((file_path, original_name, unique_name))
        
        files_list = "\n".join([f"📄 {i+1}. {name}" for i, (_, name, _) in enumerate(user_data[user_id].files)])
        
        keyboard = [
            [InlineKeyboardButton(f"🗂 Arxivlash ({len(user_data[user_id].files)})", callback_data=ZIP_BUTTON)],
            [InlineKeyboardButton("🧹 Ro'yxatni tozalash", callback_data=CLEAR_BUTTON)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        try:
            await context.bot.edit_message_text(
                chat_id=chat_id,
                message_id=user_data[user_id].status_message_id,
                text=f"📁 <b>Fayl qabul qilindi!</b>\n\n"
                     f"📂 Joriy fayllar ro'yxati: \n"
                     f"{files_list}\n\n"
                     f"ℹ️ Jami: {len(user_data[user_id].files)} ta fayl",
                reply_markup=reply_markup,
                parse_mode=ParseMode.HTML
            )
        except:
            message = await update.message.reply_text(
                f"📁 <b>Fayl qabul qilindi!</b>\n\n"
                f"📂 Joriy fayllar ro'yxati:\n"
                f"{files_list}\n\n"
                f"ℹ️ Jami: {len(user_data[user_id].files)} ta fayl",
                reply_markup=reply_markup,
                parse_mode=ParseMode.HTML
            )
            user_data[user_id].status_message_id = message.message_id
    
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik yuz berdi: {str(e)}")

async def zip_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    if user_id not in user_data or not user_data[user_id].files:
        await query.edit_message_text(text="⚠️ Arxivlash uchun fayllarni yuboring!")
        return
    
    try:
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")
        
        zip_filename = f"Arxiv_{user_id}.zip"
        
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for file_path, _, unique_name in user_data[user_id].files:
                zipf.write(file_path, unique_name)
        
        await context.bot.send_chat_action(chat_id=chat_id, action="upload_document")
        await context.bot.send_document(
            chat_id=chat_id,
            document=open(zip_filename, 'rb'),
            filename="arxiv.zip",
            caption=f"✅ {len(user_data[user_id].files)} ta fayl arxivlandi!\n 🌐 @Convert_filesbot"
        )
        for file_path, _, _ in user_data[user_id].files:
            if os.path.exists(file_path):
                os.remove(file_path)
        user_data[user_id].files = []
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
            
    except Exception as e:
        await query.edit_message_text(text=f"❌ Xatolik yuz berdi: {str(e)}")
    
    finally:
        if os.path.exists(zip_filename):
            os.remove(zip_filename)

async def clear_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = update.effective_user.id
    
    if user_id in user_data:
        for file_path, _, _ in user_data[user_id].files:
            if os.path.exists(file_path):
                os.remove(file_path)
        user_data[user_id].files = []
    
    keyboard = [
        [InlineKeyboardButton("🗂 ARXIVLASH (0)", callback_data=ZIP_BUTTON)]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await query.edit_message_text(
        text="📁 <b>Ro'yxat tozalandi!</b>\n\n"
             "🔹 Yangi fayllar yuborishingiz mumkin\n"
             "🔹 Fayllar ro'yxati bo'sh\n\n"
             "ℹ️ Jami: 0 ta fayl",
        reply_markup=reply_markup,
        parse_mode=ParseMode.HTML
    )

async def return_to_main_menu(chat_id: int, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Asosiy menyuga qaytish, avval obuna holatini tekshirish"""
    user_id = chat_id  # chat_id foydalanuvchi ID si sifatida ishlatiladi
    user = await context.bot.get_chat(user_id)
    first_name = user.first_name or "Foydalanuvchi"
    try:
        # Obuna holatini tekshirish
        is_subscribed = await check_all_subscriptions(user_id, context)
        if not is_subscribed:
            await send_subscription_request_to_chat(chat_id, context)
            return
        # Asosiy menyu tugmalari
        keyboard = [
            [InlineKeyboardButton("📄Word 🔄 PDF", callback_data='docx_pdf'),
             InlineKeyboardButton("📊Excel ➡ PDF", callback_data='excel_pdf')],
            [InlineKeyboardButton("🎤Slayd (PPTX) ➡ PDF", callback_data='ppt_pdf'),
             InlineKeyboardButton("🖼RASM(JPG) ➡ PDF", callback_data='jpg_pdf')],
            [InlineKeyboardButton("🔤Kiril 🔄 Lotin", callback_data='translate_file'),
             InlineKeyboardButton("🔲QR Kod Yasash", callback_data='qr_gen')],
            [InlineKeyboardButton("📷QR Kod Aniqlash", callback_data='qr_scan'),
             InlineKeyboardButton("💧PDF Suv belgi qo'yish", callback_data='pdf_watermark')],
            [InlineKeyboardButton("🔒PDF Parol qo'yish", callback_data='pdf_protect'),
             InlineKeyboardButton("🔓PDF Parolni olib tashlash", callback_data='pdf_unprotect')],
            [InlineKeyboardButton("📝Word ➡ Excel", callback_data='word_excel'),
             InlineKeyboardButton("🗂Arxivlovchi", callback_data='file_zipper')],
            [InlineKeyboardButton("📜Matn chiqarish", callback_data='ocr'),
             InlineKeyboardButton("🗜Fayl Siqish", callback_data='compress')],
            [InlineKeyboardButton("📄PDF Ajratish", callback_data='pdf_split'),
             InlineKeyboardButton("ℹYo'riqnoma", callback_data='about')]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
    
        welcome_text = f'''
    *Assalomu alaykum {first_name}!* 👋
O'zingizga kerakli bo'limni tanlang!
Botni qayta ishga tushurish uchun /start buyrug'ini yuboring
ℹ Botdan foydalanish yo'riqnomasi👇
        '''
        await context.bot.send_message(
            chat_id=chat_id,
            text=welcome_text,
            reply_markup=reply_markup,
            parse_mode="Markdown"
        )
    except Exception as e:
        logger.error(f"Asosiy menyuga qaytishda xato: {str(e)}")
        await context.bot.send_message(
            chat_id=chat_id,
            text="⚠️ Kutilmagan xato yuz berdi. Iltimos, /start buyrug'ini yuboring."
        )


def generate_unique_name(filename):
    name, ext = os.path.splitext(filename)
    return f"{name}_{uuid.uuid4().hex[:6]}{ext}"

def detect_script(text: str) -> str:
    cyrillic_chars = sum(1 for char in text if '\u0400' <= char <= '\u04FF')
    latin_chars = sum(1 for char in text if char.isalpha() and char.lower() in LATIN_TO_CYRILLIC)
    return 'cyrillic' if cyrillic_chars > latin_chars else 'latin'

def extract_sample_text(file_path: str, file_extension: str) -> str:
    sample_text = ""
    try:
        if file_extension == '.doc':
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(file_path)
            sample_text = doc.Content.Text[:1000]
            doc.Close()
            word.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    sample_text += para.text + " "
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            sample_text += cell.text + " "
        
        elif file_extension == '.xls':
            pythoncom.CoInitialize()
            excel = win32com.client.Dispatch("Excel.Application")
            wb = excel.Workbooks.Open(file_path)
            
            for sheet in wb.Sheets:
                used_range = sheet.UsedRange
                for row in used_range.Rows:
                    for cell in row.Columns:
                        if cell.Value and str(cell.Value).strip():
                            sample_text += str(cell.Value) + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
            
            wb.Close()
            excel.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.xlsx':
            wb = load_workbook(file_path, read_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell and str(cell).strip():
                            sample_text += str(cell) + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
        
        elif file_extension == '.ppt':
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(file_path)
            
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "TextFrame"):
                        if shape.TextFrame.HasText:
                            sample_text += shape.TextFrame.TextRange.Text + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
            
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        sample_text += shape.text + " "
                        if len(sample_text) > 1000:
                            break
                if len(sample_text) > 1000:
                    break
    
    except Exception as e:
        logger.error(f"Namuna matn olishda xatolik: {e}")
    
    return sample_text.strip()[:1000]

def translate_old_doc(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(input_path)
        
        for paragraph in doc.Paragraphs:
            paragraph.Range.Text = translate_text(paragraph.Range.Text, translation_dict, direction)
        
        new_output_path = output_path.replace('.doc', '.docx')
        doc.SaveAs2(new_output_path, FileFormat=16)
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".doc faylni tarjima qilishda xato: {e}")
        return False

def translate_old_xls(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        wb = excel.Workbooks.Open(input_path)
        
        for sheet in wb.Sheets:
            used_range = sheet.UsedRange
            for row in used_range.Rows:
                for cell in row.Columns:
                    if cell.Value and isinstance(cell.Value, str):
                        cell.Value = translate_text(cell.Value, translation_dict, direction)
        
        new_output_path = output_path.replace('.xls', '.xlsx')
        wb.SaveAs(new_output_path, FileFormat=51)
        wb.Close()
        excel.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".xls faylni tarjima qilishda xato: {e}")
        return False

def translate_old_ppt(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path)
        
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame"):
                    if shape.TextFrame.HasText:
                        shape.TextFrame.TextRange.Text = translate_text(
                            shape.TextFrame.TextRange.Text, 
                            translation_dict, 
                            direction
                        )
        
        new_output_path = output_path.replace('.ppt', '.pptx')
        presentation.SaveAs(new_output_path, 24)
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".ppt faylni tarjima qilishda xato: {e}")
        return False

def translate_file(input_path: str, output_path: str, translation_dict: dict, direction: str, file_extension: str) -> bool:
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            if file_extension == '.doc':
                docx_path = os.path.join(temp_dir, "converted.docx")
                if not convert_doc_to_docx(input_path, docx_path):
                    return False
                return translate_docx(docx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.xls':
                xlsx_path = os.path.join(temp_dir, "converted.xlsx")
                if not convert_xls_to_xlsx(input_path, xlsx_path):
                    return False
                return translate_excel(xlsx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.ppt':
                pptx_path = os.path.join(temp_dir, "converted.pptx")
                if not convert_ppt_to_pptx(input_path, pptx_path):
                    return False
                return translate_pptx(pptx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.docx':
                return translate_docx(input_path, output_path, translation_dict, direction)
            
            elif file_extension == '.xlsx':
                return translate_excel(input_path, output_path, translation_dict, direction)
            
            elif file_extension == '.pptx':
                return translate_pptx(input_path, output_path, translation_dict, direction)
            
            else:
                logger.error(f"Noto'g'ri fayl formati: {file_extension}")
                return False
                
    except Exception as e:
        logger.error(f"Fayl tarjimasida xatolik: {e}", exc_info=True)
        return False

def translate_docx(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        doc = docx.Document(input_path)
        
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text:
                    run.text = translate_text(run.text, translation_dict, direction)
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = translate_text(run.text, translation_dict, direction)
        
        doc.save(output_path)
        return True
    except Exception as e:
        logger.error(f"DOCX tarjimasida xato: {e}")
        return False

def translate_excel(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        wb = load_workbook(input_path)
        
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cell.value = translate_text(cell.value, translation_dict, direction)
        
        wb.save(output_path)
        wb.close()
        return True
    except Exception as e:
        logger.error(f"Excel tarjimasida xato: {e}")
        return False

def translate_pptx(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                run.text = translate_text(run.text, translation_dict, direction)
            
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.text:
                                            run.text = translate_text(run.text, translation_dict, direction)
        
        prs.save(output_path)
        return True
    except Exception as e:
        logger.error(f"PowerPoint tarjimasida xato: {e}")
        return False
    
def convert_doc_to_docx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(input_path)
        doc.SaveAs2(output_path, FileFormat=16)
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".doc -> .docx konvertatsiyada xato: {e}")
        return False

def convert_xls_to_xlsx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        wb = excel.Workbooks.Open(input_path)
        wb.SaveAs(output_path, FileFormat=51)
        wb.Close()
        excel.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".xls -> .xlsx konvertatsiyada xato: {e}")
        return False

def convert_ppt_to_pptx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path)
        presentation.SaveAs(output_path, 24)
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".ppt -> .pptx konvertatsiyada xato: {e}")
        return False
    
def translate_text(text: str, translation_dict: dict, direction: str) -> str:
    try:
        if not isinstance(text, str):
            return text
            
        if direction == 'lat_cyr':
            # Maxsus birikmalar uchun (katta harflar bilan boshlangan)
            text = re.sub(r'\bYe\b', 'Е', text)
            text = re.sub(r'Ye', 'Е', text)
            text = re.sub(r'\bYa\b', 'Я', text)
            text = re.sub(r'Ya', 'Я', text)
            text = re.sub(r'\bYu\b', 'Ю', text)
            text = re.sub(r'Yu', 'Ю', text)
            
            # Kichik harflar uchun
            text = re.sub(r'\bye\b', 'е', text)
            text = re.sub(r'ye', 'е', text)
            text = re.sub(r'\bya\b', 'я', text)
            text = re.sub(r'ya', 'я', text)
            text = re.sub(r'\byu\b', 'ю', text)
            text = re.sub(r'yu', 'ю', text)
            
            # Qolgan maxsus birikmalar (sh, ch, g', o')
            text = re.sub(r'\bSh\b', 'Ш', text)
            text = re.sub(r'\bSH\b', 'Ш', text)
            text = re.sub(r'\bSh', 'Ш', text)
            text = re.sub(r'SH', 'Ш', text)
            text = re.sub(r'sh', 'ш', text)
            
            text = re.sub(r'\bCh\b', 'Ч', text)
            text = re.sub(r'\bCH\b', 'Ч', text)
            text = re.sub(r'\bCh', 'Ч', text)
            text = re.sub(r'CH', 'Ч', text)
            text = re.sub(r'ch', 'ч', text)
            
            text = re.sub(r"g['‘’`]", "ғ", text)
            text = re.sub(r"G['‘’`]", "Ғ", text)
            text = re.sub(r"o['‘’`]", "ў", text)
            text = re.sub(r"O['‘’`]", "Ў", text)
            text = re.sub(r"['‘’`]", "ъ", text)
            
            # Qolgan harflar
            for lat, cyr in LATIN_TO_CYRILLIC.items():
                if len(lat) == 1 and lat not in ["'", "`"]:
                    text = text.replace(lat, cyr)
        else:
            # Kirill -> Lotin
            for cyr, lat in CYRILLIC_TO_LATIN.items():
                text = text.replace(cyr, lat)
        
        # Unicode matnni to'g'ri kodlash
        return text.encode('utf-8', 'ignore').decode('utf-8')
    
    except Exception as e:
        logger.error(f"Tarjima qilishda xato: {e}")
        return text
    
async def error_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    logger.error(msg="Exception while handling an update:", exc_info=context.error)
    if update and update.message:
        await update.message.reply_text(
            "⚠️ Kutilmagan xato yuz berdi.\n"
        "Iltimos, quyidagilarni sinab ko'ring:\n"
        "- Fayl formati to'g'ri ekanligini tekshiring\n"
        "- Fayl hajmi 10MB dan kichik ekanligiga ishonch hosil qiling\n"
        "- /start buyrug'i orqali botni qayta ishga tushiring\n"
        "Agar muammo davom etsa, @Dilxush_Bahodirov ga murojaat qiling."
        )
#############
#############
def load_channels_from_excel():
    """obuna.xlsx faylidan kanallarni yuklash yoki yangi fayl yaratish"""
    global subscription_channels
    subscription_channels = []
    
    try:
        file_path = Path("obuna.xlsx")
        if not file_path.exists():
            # Yangi Excel fayli yaratish
            wb = Workbook()
            sheet = wb.active
            sheet.append(["Kanal nomi", "Username/ID", "Link"])  # Sarlavha qatori
            wb.save(file_path)
            logger.info("Yangi obuna.xlsx fayli yaratildi")
            return
            
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        
        for row in sheet.iter_rows(min_row=2, values_only=True):
            if len(row) >= 3 and all(row[:3]):
                name, identifier, link = row[:3]
                subscription_channels.append({
                    'name': str(name).strip(),
                    'id': str(identifier).strip(),
                    'link': str(link).strip(),
                    'is_private': not str(identifier).startswith('@')
                })
        
        logger.info(f"obuna.xlsx dan {len(subscription_channels)} ta kanal yuklandi")
        
    except Exception as e:
        logger.error(f"obuna.xlsx ni o'qishda xato: {e}")

async def check_all_subscriptions(user_id: int, context: CallbackContext) -> bool:
    """Foydalanuvchi barcha kanallarga obuna bo'lganligini tekshiradi"""
    if not subscription_channels:
        return True
        
    for channel in subscription_channels:
        is_subscribed = await check_subscription(user_id, channel, context)
        if not is_subscribed:
            return False
            
    return True

async def check_subscription(user_id: int, channel_data: dict, context: CallbackContext) -> bool:
    try:
        channel_id = channel_data['id']
        
        # ID ni to'g'ri formatga keltirish
        if isinstance(channel_id, str):
            if channel_id.startswith('@'):
                # Username bilan ishlash
                pass
            elif channel_id.startswith('-100'):
                channel_id = int(channel_id)
            elif channel_id.startswith('-'):
                # -123456 -> -100123456
                channel_id = int(f"-100{channel_id[1:]}")
            else:
                # 123456 -> -100123456
                channel_id = int(f"-100{channel_id}")
        elif isinstance(channel_id, int):
            if channel_id > 0:
                channel_id = int(f"-100{channel_id}")
            elif -1000000000000 < channel_id < -999999999999:
                # Agar allaqachon -100... formatida bo'lsa
                pass
            else:
                # -123456 -> -100123456
                channel_id = int(f"-100{abs(channel_id)}")
        
        # Bot adminligini tekshirish
        try:
            chat = await context.bot.get_chat(chat_id=channel_id)
            if chat.type in ['channel', 'supergroup']:
                bot_member = await context.bot.get_chat_member(
                    chat_id=channel_id,
                    user_id=context.bot.id
                )
                if bot_member.status not in ['administrator', 'creator']:
                    logger.error(f"Bot admin emas! {channel_id}")
                    return False
        except Exception as e:
            logger.error(f"Chat topilmadi/adminlik tekshirilmadi: {e}")
            return False
        
        # Obunani tekshirish
        try:
            member = await context.bot.get_chat_member(
                chat_id=channel_id,
                user_id=user_id
            )
            return member.status in ['member', 'administrator', 'creator']
        except Exception as e:
            logger.error(f"Obuna tekshirishda xato: {e}")
            return False
            
    except Exception as e:
        logger.error(f"Umumiy xato: {e}")
        return False
    
async def update_subscription_list(user_id: int, context: CallbackContext):
    """Obuna bo'lgan kanallarni ro'yxatdan olib tashlash"""
    global subscription_channels
    
    # Yangi ro'yxat - faqat obuna bo'lmagan kanallar
    new_channels = []
    
    for channel in subscription_channels:
        # channel bu dictionary bo'lishi kerak {'id': ..., 'name': ...}
        if not isinstance(channel, dict):
            logger.error(f"Noto'g'ri kanal formati: {channel}")
            continue
            
        is_subscribed = await check_subscription(user_id, channel, context)
        if not is_subscribed:
            new_channels.append(channel)
    
    # Yangi ro'yxatni o'rnatish
    subscription_channels = new_channels

async def send_subscription_request_to_chat(chat_id: int, context: ContextTypes.DEFAULT_TYPE):
    """Foydalanuvchiga faqat obuna bo'lmagan kanallar haqida xabar yuborish"""
    user_id = chat_id
    try:
        # Avvalgi xabarni o'chirish
        if user_id in pending_messages:
            try:
                await context.bot.delete_message(
                    chat_id=user_id,
                    message_id=pending_messages[user_id]
                )
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")

        # Foydalanuvchi obuna bo'lmagan kanallarni aniqlash
        unsubscribed_channels = []
        for channel in subscription_channels:
            is_subscribed = await check_subscription(user_id, channel, context)
            if not is_subscribed:
                unsubscribed_channels.append(channel)

        # Agar barcha kanallarga obuna bo'lsa
        if not unsubscribed_channels:
            await send_confirmation_message(user_id, context)
            return

        # Obuna bo'lish kerak bo'lgan kanallar uchun tugmalar yaratish
        keyboard = []
        for channel in unsubscribed_channels:
            keyboard.append([InlineKeyboardButton(
                f"Obuna bo'lish: {channel['name']}", 
                url=channel['link']
            )])

        reply_markup = InlineKeyboardMarkup(keyboard)
        
        message_text = "Botdan foydalanish uchun quyidagi kanallarga obuna bo'ling:\n\n"
        message_text += "\n".join([f"🔹 {channel['name']}" for channel in unsubscribed_channels])
        message_text += "\n\nObuna bo'lgach, istalgan xabarni yuboring."

        message = await context.bot.send_message(
            chat_id=user_id,
            text=message_text,
            reply_markup=reply_markup
        )
        pending_messages[user_id] = message.message_id
        
        # Obunani avtomatik tekshirishni boshlash
        await start_subscription_check(user_id, context)
    except Exception as e:
        logger.error(f"Obuna so'rovi yuborishda xato: {str(e)}")

async def send_subscription_request(update: Update, context: CallbackContext):
    """Foydalanuvchiga faqat obuna bo'lmagan kanallar haqida xabar yuboradi"""
    user_id = update.effective_user.id
    
    # Avvalgi xabarni o'chirish
    if user_id in pending_messages:
        try:
            await context.bot.delete_message(
                chat_id=user_id,
                message_id=pending_messages[user_id]
            )
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")

    # Foydalanuvchi obuna bo'lmagan kanallarni aniqlash
    unsubscribed_channels = []
    for channel in subscription_channels:
        is_subscribed = await check_subscription(user_id, channel, context)
        if not is_subscribed:
            unsubscribed_channels.append(channel)

    # Agar barcha kanallarga obuna bo'lsa
    if not unsubscribed_channels:
        await send_confirmation_message(user_id, context)
        return

    # Obuna bo'lish kerak bo'lgan kanallar uchun tugmalar yaratish
    keyboard = []
    for channel in unsubscribed_channels:
        keyboard.append([InlineKeyboardButton(
            f"Obuna bo'lish: {channel['name']}", 
            url=channel['link']
        )])

    reply_markup = InlineKeyboardMarkup(keyboard)
    
    message_text = "Botdan foydalanish uchun quyidagi kanallarga obuna bo'ling:\n\n"
    message_text += "\n".join([f"🔹 {channel['name']}" for channel in unsubscribed_channels])
    message_text += "\n\nObuna bo'lgach, istalgan xabarni yuboring."

    message = await update.message.reply_text(
        message_text,
        reply_markup=reply_markup
    )
    pending_messages[user_id] = message.message_id
    
    # Obunani avtomatik tekshirishni boshlash
    await start_subscription_check(user_id, context)

async def start_subscription_check(user_id: int, context: CallbackContext):
    """Obunani avtomatik tekshirishni boshlash"""
    global check_tasks
    
    if user_id in check_tasks:
        check_tasks[user_id].cancel()
    
    async def check_loop():
        while True:
            await asyncio.sleep(5)
            
            remaining_channels = []
            for channel in subscription_channels:
                is_subscribed = await check_subscription(user_id, channel, context)
                if not is_subscribed:
                    remaining_channels.append(channel)
            
            if not remaining_channels:
                if user_id in pending_messages:
                    try:
                        await context.bot.delete_message(
                            chat_id=user_id,
                            message_id=pending_messages[user_id]
                        )
                        del pending_messages[user_id]
                    except Exception as e:
                        logger.error(f"Xabarni o'chirishda xato: {e}")
                
                await send_confirmation_message(user_id, context)
                
                if user_id in check_tasks:
                    del check_tasks[user_id]
                break
            
            elif user_id in pending_messages:
                try:
                    keyboard = []
                    for channel in remaining_channels:
                        keyboard.append([InlineKeyboardButton(
                            f"Obuna bo'lish: {channel['name']}", 
                            url=channel['link']
                        )])
                    
                    reply_markup = InlineKeyboardMarkup(keyboard)
                    
                    new_text = "Botdan foydalanish uchun quyidagi kanallarga obuna bo'ling:\n\n"
                    new_text += "\n".join([f"🔹 {channel['name']}" for channel in remaining_channels])
                    new_text += "\n\nObuna bo'lgach, istalgan xabarni yuboring."
                    
                    await context.bot.edit_message_text(
                        chat_id=user_id,
                        message_id=pending_messages[user_id],
                        text=new_text,
                        reply_markup=reply_markup
                    )
                except Exception as e:
                    logger.error(f"Xabarni yangilashda xato: {e}")
    
    task = asyncio.create_task(check_loop())
    check_tasks[user_id] = task

async def generate_users_csv():
    """Foydalanuvchilar ro'yxatini CSV formatida yaratish"""
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(['User ID', 'First Name', 'Username'])
    
    for user_id, user_info in users.items():
        writer.writerow([
            user_id,
            user_info.get('first_name', 'N/A'),
            user_info.get('username', 'N/A')
        ])
    
    output.seek(0)
    return output

async def add_channel(update: Update, context: CallbackContext):
    """Kanal qo'shish - yangi ID qoidalari bilan"""
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari mavjud emas!")
        return
    
    try:
        # Xabarni 'Nomi|ID/@username|Link' formatida ajratamiz
        parts = update.message.text.split('|', 2)
        if len(parts) != 3:
            await update.message.reply_text(
                "❌ Format xato! Quyidagi ko'rinishda yuboring:\n"
                "Nomi|ID/@username|Link\n\n"
                "Misol:\n"
                "Test Kanal|@testchannel|https://t.me/testchannel\n"
                "Yopiq Guruh|-2383562706|https://t.me/c/2383562706/1\n"
                "Boshqa Guruh|2383562706|https://t.me/c/2383562706/1"
            )
            return
            
        name, identifier, link = parts
        name = name.strip()
        identifier = identifier.strip()
        link = link.strip()
        
        # Kanal allaqachon mavjudligini qat'iy tekshiramiz
        for channel in subscription_channels:
            # ID yoki username bir xil bo'lsa
            if channel['id'].lower() == identifier.lower():
                await update.message.reply_text(
                    f"❌ Bu kanal ({identifier}) allaqachon ro'yxatda mavjud!\n"
                    f"Kanal nomi: {channel['name']}\n"
                    f"Link: {channel['link']}"
                )
                return
            
            # Link bir xil bo'lsa
            if channel['link'].lower() == link.lower():
                await update.message.reply_text(
                    f"❌ Bu link allaqachon boshqa kanal uchun ishlatilgan!\n"
                    f"Mavjud kanal: {channel['name']} ({channel['id']})\n"
                    f"Link: {channel['link']}"
                )
                return
            
            # Nomi bir xil bo'lsa (ixtiyoriy, agar kerak bo'lsa)
            if channel['name'].lower() == name.lower():
                await update.message.reply_text(
                    f"⚠️ Diqqat: {name} nomli kanal allaqachon mavjud!\n"
                    f"Mavjud kanal ID: {channel['id']}\n"
                    f"Yangi kiritilgan ID: {identifier}\n\n"
                    "Agar bu boshqa kanal bo'lsa, nomini o'zgartiring."
                )
                return
            
        # Identifikator turini aniqlaymiz
        is_private = False
        if identifier.startswith('@'):
            # Username (@ bilan boshlanadi)
            #if not identifier[1:].isalnum():  # Faqat harf va raqamlar
                #await update.message.reply_text("❌ Noto'g'ri username format! Faqat harflar, raqamlar va _ belgisi bo'lishi mumkin")
                #return
            if not (link.startswith('t.me/') or link.startswith('https://t.me/')):
                await update.message.reply_text("❌ Noto'g'ri kanal link formati! t.me/ yoki https://t.me/ bilan boshlanishi kerak")
                return
        elif identifier.startswith('-'):
            # ID (- bilan boshlanadi)
            if not identifier[1:].isdigit():
                await update.message.reply_text("❌ Noto'g'ri ID formati! - dan keyin faqat raqamlar bo'lishi kerak")
                return
            is_private = True
            if not (link.startswith('t.me/') or link.startswith('https://t.me/')):
                await update.message.reply_text("❌ Noto'g'ri guruh link formati! t.me/+ yoki https://t.me/+ bilan boshlanishi kerak")
                return
        else:
            # ID (oddiy raqam)
            if not identifier.isdigit():
                await update.message.reply_text("❌ Noto'g'ri ID formati! Faqat raqamlar bo'lishi kerak")
                return
            if not (link.startswith('t.me/') or link.startswith('https://t.me/')):
                await update.message.reply_text("❌ Noto'g'ri guruh link formati! t.me/+ yoki https://t.me/+ bilan boshlanishi kerak")
                return
        
        # Kanal allaqachon mavjudligini tekshiramiz
        for channel in subscription_channels:
            if channel['id'].lower() == identifier.lower():
                await update.message.reply_text(f"❌ Bu chat ({identifier}) allaqachon ro'yxatda mavjud!")
                return
        
        # Yangi kanal ma'lumotlarini yaratamiz
        new_channel = {
            'name': name,
            'id': identifier,
            'link': link if link.startswith('http') else f'https://{link}',
            'is_private': not identifier.startswith('@')
        }
        
        # Kanalni ro'yxatga qo'shamiz
        subscription_channels.append(new_channel)
        
        # Excel fayliga yozamiz
        save_channels_to_excel()
        
        await update.message.reply_text(
            f"✅ {'Guruh' if new_channel['is_private'] else 'Kanal'} muvaffaqiyatli qo'shildi:\n"
            f"📌 Nomi: {name}\n"
            f"🆔 {'ID' if new_channel['is_private'] else 'Username'}: {identifier}\n"
            f"🔗 Link: {link}\n\n"
            f"Jami obuna kanallari: {len(subscription_channels)}"
        )
        
    except Exception as e:
        logger.error(f"Kanal qo'shishda xato: {e}")
        await update.message.reply_text("❌ Kanal qo'shishda xatolik yuz berdi. Iltimos, qayta urunib ko'ring.")
        
def save_channels_to_excel():
    """Kanalar ro'yxatini obuna.xlsx fayliga saqlash"""
    try:
        file_path = Path("obuna.xlsx")
        
        # Agar fayl mavjud bo'lmasa, yangi yaratamiz
        if not file_path.exists():
            wb = Workbook()
            sheet = wb.active
            sheet.append(["Kanal nomi", "Username/ID", "Link"])  # Sarlavha qatorini qo'shamiz
        else:
            wb = openpyxl.load_workbook(file_path)
            sheet = wb.active
        
        # Hozirgi ma'lumotlarni tozalash (sarlavhadan tashqari)
        sheet.delete_rows(2, sheet.max_row)  # 2-qatordan boshlab hammasini o'chiramiz
        
        # Yangi ma'lumotlarni yozish
        for channel in subscription_channels:
            sheet.append([channel['name'], channel['id'], channel['link']])
        
        wb.save(file_path)
        logger.info("obuna.xlsx fayliga kanallar muvaffaqiyatli saqlandi")
        
    except Exception as e:
        logger.error(f"obuna.xlsx fayliga yozishda xato: {e}")


async def remove_channel(update: Update, context: CallbackContext):
    """Kanalni olib tashlash va Excel faylini yangilash"""
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari mavjud emas!")
        return
    
    try:
        if not context.args:
            # Kanal ro'yxatini ko'rsatish
            if not subscription_channels:
                await update.message.reply_text("❌ Hozircha kanallar qo'shilmagan.")
                return
            
            message = "Olib tashlash uchun kanal raqamini tanlang:\n\n"
            for i, channel in enumerate(subscription_channels, 1):
                message += f"{i}. {channel['name']} ({channel['id']})\n"
            
            message += "\nFoydalanish: /remove <kanal_raqami>"
            await update.message.reply_text(message)
            return
            
        channel_num = int(context.args[0])
        
        # Raqamni tekshirish (1 dan len(subscription_channels) gacha bo'lishi kerak)
        if 1 <= channel_num <= len(subscription_channels):
            removed_channel = subscription_channels.pop(channel_num - 1)  # Ro'yxat indeksi 0 dan boshlanadi
            
            # Excel faylini yangilaymiz
            save_channels_to_excel()
            
            await update.message.reply_text(
                f"✅ Kanal muvaffaqiyatli o'chirildi:\n"
                f"📌 Nomi: {removed_channel['name']}\n"
                f"🆔 ID: {removed_channel['id']}\n"
                f"🔗 Link: {removed_channel['link']}\n\n"
                f"Qolgan kanallar soni: {len(subscription_channels)}"
            )
        else:
            await update.message.reply_text(
                f"❌ Noto'g'ri kanal raqami! 1 dan {len(subscription_channels)} gacha raqam kiriting.\n"
                f"Jami kanallar: {len(subscription_channels)}"
            )
    except ValueError:
        await update.message.reply_text("❌ Noto'g'ri raqam formati! Faqat raqam kiriting.")
    except Exception as e:
        logger.error(f"Kanalni olib tashlashda xato: {e}")
        await update.message.reply_text("❌ Kanalni olib tashlashda xatolik yuz berdi.")

async def list_channels(update: Update, context: CallbackContext):
    """Kanalar ro'yxatini ko'rsatish"""
    if update.effective_user.id != ADMIN_ID:
        return
    
    if not subscription_channels:
        await update.message.reply_text("❌ Hozircha kanallar qo'shilmagan.")
        return
    
    message = "📋 Majburiy obuna kanallari:\n\n"
    for i, channel in enumerate(subscription_channels, 1):
        message += f"{i}. {channel['name']} ({channel['id']})\n"
    
    await update.message.reply_text(message)

async def show_statistics(update: Update, context: CallbackContext):
    """Bot statistikasini ko'rsatish"""
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari mavjud emas!")
        return
    
    try:
        stats = (
            f"📊 Bot statistikasi:\n\n"
            f"👥 Foydalanuvchilar soni: {len(users)}\n"
            f"🔗 Majburiy obuna kanallari: {len(subscription_channels)}\n"
            f"📅 Oxirgi faollik: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        )
        
        if subscription_channels:
            stats += "\n\n📋 Kanalar ro'yxati:\n"
            for i, channel in enumerate(subscription_channels, 1):
                stats += f"{i}. {channel['name']} ({channel['id']})\n"
        
        await update.message.reply_text(stats)
    except Exception as e:
        logger.error(f"Statistika ko'rsatishda xato: {e}")
        await update.message.reply_text("❌ Statistika yuklashda xatolik yuz berdi.")
   
async def send_confirmation_message(user_id: int, context: ContextTypes.DEFAULT_TYPE):
    
    # Avvalgi tasdiqlash xabarini o'chirish
    if user_id in confirmation_messages:
        try:
            await context.bot.delete_message(
                chat_id=user_id,
                message_id=confirmation_messages[user_id]
            )
        except Exception as e:
            logger.error(f"Tasdiqlash xabarini o'chirishda xato: {e}")
    
    # Yangi tasdiqlash xabarini yuborish
    message = await context.bot.send_message(
        chat_id=user_id,
        text="✅ Obunangiz tasdiqlandi! Endi botdan to'liq foydalana olasiz.\n\n"
             "Botdan foydalanish uchun /start buyrug'ini yuboring."
    )
    confirmation_messages[user_id] = message.message_id
    
    # 10 soniyadan keyin xabarni o'chirish va asosiy menyuga qaytish
    await asyncio.sleep(5)
    try:
        await context.bot.delete_message(
            chat_id=user_id,
            message_id=message.message_id
        )
        if user_id in confirmation_messages:
            del confirmation_messages[user_id]
        
        # Asosiy menyuga qaytish
        await return_to_main_menu(user_id, context)
    except Exception as e:
        logger.error(f"Tasdiqlash xabarini o'chirishda xato: {e}")


async def delete_confirmation_message(context: CallbackContext):
    """Tasdiqlash xabarini 5 soniyadan keyin o'chiradi"""
    job = context.job
    try:
        await context.bot.delete_message(
            chat_id=job.chat_id,
            message_id=job.data
        )
    except Exception as e:
        logger.error(f"Xabarni o'chirishda xato: {e}")

async def handle_message(update: Update, context: CallbackContext) -> None:
    """Barcha xabarlarni qayta ishlash"""
    global messages_sent
    user_id = update.effective_user.id
    text = update.message.text
    
    # Admin xabarlarini alohida ishlash
    if user_id == ADMIN_ID:
        await handle_admin_message(update, context)
        return
    
    # Obunani tekshirish
    if subscription_channel:
        # Agar avval obuna bo'lgan bo'lsa
        if user_subscriptions.get(user_id, False):
            await update.message.reply_text("Xabaringiz qabul qilindi!")
            messages_sent += 1
            return
        
        # Obunani tekshiramiz
        is_subscribed = await check_subscription(user_id, context)
        
        if is_subscribed:
            user_subscriptions[user_id] = True
            
            # "Obuna bo'ling" xabarini o'chirish
            if user_id in pending_messages:
                try:
                    await context.bot.delete_message(
                        chat_id=user_id,
                        message_id=pending_messages[user_id]
                    )
                    del pending_messages[user_id]
                except Exception as e:
                    logger.error(f"Xabarni o'chirishda xato: {e}")
            
            # Yangi tasdiqlash xabarini yuborish
            confirmation_msg = await update.message.reply_text(
                "✅ Obunangiz tasdiqlandi! Endi botdan to'liq foydalana olasiz.\n\n"
                "Botdan foydalanish uchun /start buyrug'ini yuboring."
            )
            
            # 5 soniyadan keyin tasdiqlash xabarini o'chirish
            context.job_queue.run_once(
                delete_confirmation_message,
                when=timedelta(seconds=5),
                chat_id=user_id,
                data=confirmation_msg.message_id,
                name=f"del_conf_{confirmation_msg.message_id}"
            )
            
            messages_sent += 1
           
        else:
            # Agar obuna bo'lmagan bo'lsa, yangi obuna so'rovini yuboramiz
            if user_id not in pending_messages:  # Faqat bir marta yuborish uchun
                await send_subscription_request(update, context)
        
    else:
        await update.message.reply_text("Xabaringiz qabul qilindi!")
        messages_sent += 1
        
async def admin(update: Update, context: CallbackContext) -> None:
    """Admin komandasi handleri"""
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari mavjud emas!")
        return
    
    await admin_panel(update, context)
    
async def handle_message(update: Update, context: CallbackContext) -> None:
    """Barcha xabarlarni qayta ishlash"""
    user_id = update.effective_user.id
    user = update.effective_user
    users[user_id] = {
        'first_name': user.first_name or 'N/A',
        'username': f"@{user.username}" if user.username else 'N/A'
    }
    
    if user_id == ADMIN_ID:
        await handle_admin_message(update, context)
        return
    
    if subscription_channels:
        is_subscribed = await check_all_subscriptions(user_id, context)
        if not is_subscribed:
            if user_id not in pending_messages:
                await send_subscription_request(update, context)
            return
    
    await update.message.reply_text("Xabaringiz qabul qilindi!")

async def handle_admin_message(update: Update, context: CallbackContext):
    """Admin panelidagi harakatlarni boshqarish"""
    user_id = update.effective_user.id
    if user_id != ADMIN_ID:
        return
    
    text = update.message.text
    
    if text == "📋 Kanalar ro'yxati":
        if not subscription_channels:
            await update.message.reply_text("❌ Hozircha kanallar qo'shilmagan.")
            return
        
        message = "📋 Majburiy obuna kanallari:\n\n"
        for i, channel in enumerate(subscription_channels, 1):
            message += f"{i}. {channel['name']} ({channel['id']})\n"
        
        await update.message.reply_text(message)
    
    elif text == "📊 Statistika":
        await show_statistics(update, context)
    
    elif text == "➕ Kanal qo'shish":
        await update.message.reply_text(
            "Yangi kanal qo'shish uchun quyidagi formatda yuboring:\n"
            "Nomi|@username|Link\n\n"
            "Masalan: Test Kanal|@testchannel|https://t.me/testchannel"
        )
        context.user_data['awaiting_channel'] = True
    
    elif text == "➖ Kanal olib tashlash":
        if not subscription_channels:
            await update.message.reply_text("❌ Hozircha kanallar qo'shilmagan.")
            return
        
        message = "Olib tashlash uchun kanal raqamini tanlang:\n\n"
        for i, channel in enumerate(subscription_channels, 1):
            message += f"{i}. {channel['name']} ({channel['id']})\n"
        
        message += "\n/remove <kanal_raqami> buyrug'ini yuboring"
        await update.message.reply_text(message)
    
    elif text == "📢 Xabar yuborish":
        await update.message.reply_text("Barcha foydalanuvchilarga yuboriladigan xabarni kiriting:")
        context.user_data['sending_message'] = True
    
    elif text == "👥 Foydalanuvchilar ro'yxati":
        csv_file = await generate_users_csv()
        csv_bytes = io.BytesIO(csv_file.getvalue().encode('utf-8'))
        csv_bytes.name = 'users_list.csv'
        await update.message.reply_document(
            document=csv_bytes,
            filename='users_list.csv',
            caption="📋 Foydalanuvchilar ro'yxati"
        )
        csv_file.close()
    
    elif text == "🔙 Chiqish":
        await update.message.reply_text(
            "Admin panel yopildi.",
            reply_markup=ReplyKeyboardMarkup([[KeyboardButton("/start")]], resize_keyboard=True)
        )
    
    elif context.user_data.get('sending_message', False):
        await send_broadcast(context.bot, text)
        await update.message.reply_text("✅ Xabar barcha foydalanuvchilarga yuborildi!")
        del context.user_data['sending_message']
    
    elif context.user_data.get('awaiting_channel', False):
        await add_channel(update, context)
        del context.user_data['awaiting_channel']

async def setup_subscription(update: Update, context: CallbackContext) -> None:
    """Majburiy obuna sozlash"""
    await update.message.reply_text(
        "Majburiy obuna kanalini sozlang:\n\n"
        "Kanal nomi, username va linkini quyidagi formatda yuboring:\n"
        "Nomi|@username|Link\n\n"
        "Masalan: Test Kanal|@testchannel|https://t.me/testchannel"
    )
    context.user_data['awaiting_channel'] = True

async def handle_channel_input(update: Update, context: CallbackContext) -> None:
    """Kanal ma'lumotlarini qabul qilish"""
    try:
        name, username, link = update.message.text.split('|', 2)
        name = name.strip()
        username = username.strip()
        link = link.strip()
        
        if not username.startswith('@'):
            await update.message.reply_text("❌ Kanal username @ bilan boshlanishi kerak!")
            return
        
        if not (link.startswith('t.me/') or link.startswith('https://t.me/')):
            await update.message.reply_text("❌ Noto'g'ri link formati!")
            return
        
        global subscription_channel
        subscription_channel = {
            'name': name,
            'id': username,
            'link': link if link.startswith('http') else f'https://{link}'
        }
        
        await update.message.reply_text(
            f"✅ Majburiy obuna sozlandi:\n"
            f"Kanal nomi: {name}\n"
            f"Username: {username}\n"
            f"Link: {link}"
        )
        del context.user_data['awaiting_channel']
        await admin_panel(update, context)
    except ValueError:
        await update.message.reply_text("❌ Format xato! 'Nomi|@username|Link' shaklida yuboring.")
        
async def admin_panel(update: Update, context: CallbackContext):
    """Admin paneli"""
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari mavjud emas!")
        return
    
    buttons = [
        [KeyboardButton("📋 Kanalar ro'yxati"), KeyboardButton("📊 Statistika")],
        [KeyboardButton("➕ Kanal qo'shish"), KeyboardButton("➖ Kanal olib tashlash")],
        [KeyboardButton("📢 Xabar yuborish"), KeyboardButton("👥 Foydalanuvchilar ro'yxati")],
        [KeyboardButton("🔙 Chiqish")]
    ]
    reply_markup = ReplyKeyboardMarkup(buttons, resize_keyboard=True)
    await update.message.reply_text('Admin panel:', reply_markup=reply_markup)

async def send_broadcast(bot, message):
    """Barcha foydalanuvchilarga xabar yuborish"""
    success = 0
    failed = 0
    for user_id in users:
        try:
            await bot.send_message(chat_id=user_id, text=message)
            success += 1
        except Exception as e:
            logger.error(f"Xabar yuborishda xato (user {user_id}): {e}")
            failed += 1
    return success, failed

#############
#############
async def clear_pending(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Faqat admin!")
        return
    pending_messages.clear()
    await update.message.reply_text("✅ Kutilayotgan xabarlar tozalandi!")
    
    # Agar boshqa tugmalar bo'lsa, shu yerda qo'shing (masalan, docx_pdf uchun handle_docx_pdf ni alohida handler sifatida qo'shing)# Webhook Flask app
flask_app = Flask(__name__)

application = Application.builder().token(TOKEN).build()

def main() -> None:
    """Lokal polling uchun asosiy funksiya"""
    load_channels_from_excel()  # Agar allaqachon chaqirilgan bo'lmasa
    
    # Handler'larni qo'shish
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("admin", admin))
    application.add_handler(CommandHandler("addchannel", add_channel))
    application.add_handler(CommandHandler("remove", remove_channel))
    application.add_handler(CommandHandler("stats", show_statistics))
    application.add_handler(CommandHandler("clearpending", clear_pending))
    application.add_handler(CallbackQueryHandler(button_handler))
    
    # Message handler'lar
    docx_pdf_handler = MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_docx_pdf)
    translation_handler = MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_translation_doc)
    zip_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO | filters.VIDEO | filters.AUDIO) & ~filters.COMMAND, 
        handle_zip_file
    )
    ocr_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO) & ~filters.COMMAND, 
        handle_ocr
    )
    compress_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO) & ~filters.COMMAND, 
        handle_compress
    )
    
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    application.add_handler(docx_pdf_handler, group=1)
    application.add_handler(translation_handler, group=2)
    application.add_handler(zip_handler, group=3)
    application.add_handler(ocr_handler, group=4)
    application.add_handler(compress_handler, group=5)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_excel_pdf), group=6)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_ppt_pdf), group=7)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_qr_gen), group=8)
    application.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, handle_qr_scan), group=9)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_watermark), group=10)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_watermark_text), group=10)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_protect), group=11)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_protect_password), group=11)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_unprotect), group=12)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_unprotect_password), group=12)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_word_excel), group=13)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_split), group=14)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_split_pages), group=14)
    
    application.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, handle_photo))
    application.add_error_handler(error_handler)
    
    # Polling ishga tushirish
    application.run_polling(drop_pending_updates=True)  # drop_pending_updates=True - eski xabarlarni e'tiborsiz qoldirish

@flask_app.route('/webhook', methods=['POST'])
def webhook():
    try:
        json_data = request.get_json()
        update = Update.de_json(json_data, application.bot)
        asyncio.run(application.process_update(update))
        return 'OK'
    except Exception as e:
        logger.error(f"Webhook xatosi: {e}")
        return 'Error', 500

@flask_app.route('/')
def index():
    return "Telegram Bot ishga tushdi! Webhook: /webhook"

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    # Webhook o'rnatish (birinchi ishga tushganda)
    webhook_url = f"https://{os.environ.get('RENDER_EXTERNAL_HOSTNAME', 'localhost:5000')}/webhook"
    asyncio.run(application.bot.set_webhook(url=webhook_url))
    logger.info(f"Webhook o'rnatildi: {webhook_url}")
    flask_app.run(host='0.0.0.0', port=port)
